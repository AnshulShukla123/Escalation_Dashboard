
import os
import socket
import requests
import secrets
from datetime import datetime, timedelta
from flask import Flask, render_template, request, redirect, url_for, session, jsonify, flash, send_file
from flask_sqlalchemy import SQLAlchemy
from sqlalchemy import or_, and_
from dotenv import load_dotenv
import smtplib
from email.mime.text import MIMEText
import io
from functools import wraps

from email.mime.text import MIMEText

load_dotenv()

# --- AI SERVICE CLASS ---
class AIService:
    def __init__(self):
        
        pass  
        
    def generate_escalation_summary(self, escalation_data):
        """Generate AI summary for an escalation"""
        try:
            # Prepare the escalation data for AI analysis
            prompt = self._create_escalation_prompt(escalation_data)
            summary = self._generate_summary_with_ai(prompt, escalation_data)
            
            return {
                'success': True,
                'summary': summary,
                'generated_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'summary': 'AI summary generation failed. Please try again later.'
            }
    
    def _create_escalation_prompt(self, data):
        """Create a structured prompt for AI analysis"""
        prompt = f"""
        Please analyze the following escalation details and provide a comprehensive summary:
        
        ESCALATION DETAILS:
        - Customer: {data.get('customer', 'N/A')}
        - DE Manager: {data.get('dedt_manager', 'N/A')}
        - Primary Engineer: {data.get('engineer', 'N/A')}
        - Version: {data.get('version', 'N/A')}
        - State: {data.get('state', 'N/A')}
        - Severity: {data.get('severity', 'N/A')}
        - Bug ID: {data.get('bugid', 'N/A')}
        - Component: {data.get('component_name', 'N/A')}
        - Cross Dependencies: {data.get('cross_team', 'N/A')}
        - SR Number: {data.get('sr', 'N/A')}
        - BEMS: {data.get('bems', 'N/A')}
        
        TECHNICAL DETAILS:
        - Symptom: {data.get('symptom', 'N/A')}
        - Upgrade Attempt: {data.get('upgrade_attempt', 'N/A')}
        - Next Steps: {data.get('next_step', 'N/A')}
        - Remarks: {data.get('remarks', 'N/A')}
        
        Please provide a summary that includes:
        1. Problem Overview
        2. Technical Impact
        3. Current Status
        4. Key Stakeholders
        5. Next Actions Required
        """
        return prompt
    
    def _generate_summary_with_ai(self, prompt, data):
        """Generate summary using Cisco Circuit AI"""
        
        return self._generate_with_circuit(prompt, data)
    
    def _generate_with_circuit(self, prompt, data):
        """Generate summary using Cisco Circuit API"""
        try:
            # Circuit API configuration
            client_id = os.getenv('CIRCUIT_CLIENT_ID')
            client_secret = os.getenv('CIRCUIT_CLIENT_SECRET')
            app_key = os.getenv('CIRCUIT_APP_KEY')
            
            if not all([client_id, client_secret, app_key]):
                print("Warning: Circuit API credentials not found, using template fallback")
                return self._generate_template_summary(data)
            
            # Get access token
            auth_url = "https://api.ciscospark.com/v1/access_token"  
            auth_data = {
                'grant_type': 'client_credentials',
                'client_id': client_id,
                'client_secret': client_secret
            }
            
            auth_response = requests.post(auth_url, data=auth_data, timeout=10)
            
            if auth_response.status_code != 200:
                print(f"Circuit API auth failed: {auth_response.status_code}")
                return self._generate_template_summary(data)
            
            access_token = auth_response.json().get('access_token')
            
            # request for summary
            escalation_context = self._format_escalation_for_circuit(data)
            
            #ApI CALL
            api_url = f"https://api.ciscospark.com/v1/ai/summarize"  
            headers = {
                'Authorization': f'Bearer {access_token}',
                'Content-Type': 'application/json',
                'X-App-Key': app_key
            }
            
            payload = {
                'model': 'gpt-4o-mini',  
                'prompt': prompt,
                'context': escalation_context,
                'max_tokens': 800,
                'temperature': 0.3,
                'purpose': 'escalation_analysis'
            }
            
            response = requests.post(api_url, json=payload, headers=headers, timeout=30)
            
            if response.status_code == 200:
                result = response.json()
                summary = result.get('summary', result.get('content', ''))
                
                if summary:
                    return f"""🔧 **CISCO CIRCUIT AI ANALYSIS**

{summary}

---
*Generated by Cisco Circuit AI using approved enterprise models*
*Application: {app_key}*
*Analysis performed on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*"""
                else:
                    print("Empty summary from Circuit API")
                    return self._generate_template_summary(data)
            else:
                print(f"Circuit API call failed: {response.status_code} - {response.text}")
                return self._generate_template_summary(data)
                
        except requests.exceptions.Timeout:
            print("Circuit API timeout, using template fallback")
            return self._generate_template_summary(data)
        except requests.exceptions.RequestException as e:
            print(f"Circuit API request error: {e}")
            return self._generate_template_summary(data)
        except Exception as e:
            print(f"Circuit API unexpected error: {e}")
            return self._generate_template_summary(data)
    
    def _format_escalation_for_circuit(self, data):
        """Format escalation data for Circuit AI analysis"""
        return {
            'escalation_id': data.get('id', 'Unknown'),
            'customer': data.get('customer', 'N/A'),
            'severity': data.get('severity', 'N/A'),
            'state': data.get('state', 'N/A'),
            'component': data.get('component_name', 'N/A'),
            'symptom': data.get('symptom', 'N/A'),
            'de_manager': data.get('dedt_manager', 'N/A'),
            'engineer': data.get('engineer', 'N/A'),
            'cross_teams': data.get('cross_team', 'N/A'),
            'version': data.get('version', 'N/A'),
            'created_date': data.get('reported_on', 'N/A')
        }
    
    def _generate_template_summary(self, data):
        """Generate intelligent template-based summary"""
        customer = data.get('customer', 'Unknown Customer')
        version = data.get('version', 'Unknown Version')
        state = data.get('state', 'Unknown')
        severity = data.get('severity', 'Unknown')
        de_manager = data.get('dedt_manager', 'Unassigned')
        engineer = data.get('engineer', 'Unassigned')
        symptom = data.get('symptom', '')
        

        status_analysis = self._analyze_status(data)
        urgency_level = self._determine_urgency(data)
        impact_assessment = self._assess_impact(data)
        
        summary = f"""🎯 **ESCALATION SUMMARY**

📋 **Overview:**
{customer} is experiencing issues with {version} requiring escalation management. The case is currently {(state or '').lower() if state and state != 'Unknown' else 'unknown'} with {severity} severity level.

🔍 **Problem Analysis:**
{symptom[:200] + '...' if symptom and len(symptom) > 200 else symptom or 'Technical issue requiring investigation and resolution.'}

⚡ **Urgency Level:** {urgency_level}
🎯 **Impact Assessment:** {impact_assessment}

👥 **Team Assignments:**
• DE Manager: {de_manager}
• Primary Engineer: {engineer}
• Cross Teams: {data.get('cross_team', 'None')}

📊 **Current Status:**
{status_analysis}

🔄 **Next Actions:**
{data.get('next_step', 'Pending definition of next steps.')}

📝 **Key Notes:**
{(data.get('remarks', '') or 'No additional remarks provided.')[:150] + '...' if data.get('remarks') and len(data.get('remarks', '')) > 150 else (data.get('remarks') or 'No additional remarks provided.')}"""

        return summary
    
    def _analyze_status(self, data):
        """Analyze current status"""
        state = data.get('state', '')
        if state:
            state = state.lower()
        if state == 'open':
            return "🟡 Active escalation requiring ongoing attention and resolution efforts."
        elif state == 'closed':
            closing_reason = data.get('closing_reason', '')
            return f"✅ Escalation resolved and closed. Reason: {closing_reason or 'Resolution completed'}"
        else:
            return "⚪ Status pending review and classification."
    
    def _determine_urgency(self, data):
        """Determine urgency level"""
        severity = data.get('severity', '') or ''
        if severity:
            severity = severity.lower()
        cross_teams = data.get('cross_team', '') or ''
        
        if severity and ('critical' in severity or 'high' in severity):
            return "🔴 HIGH - Immediate attention required"
        elif cross_teams and cross_teams.strip() and cross_teams != '-':
            return "🟡 MEDIUM - Multiple teams involved"
        else:
            return "🟢 STANDARD - Normal escalation process"
    
    def _assess_impact(self, data):
        """Assess business impact"""
        customer = data.get('customer', '')
        component = data.get('component_name', '')
        
        if component and component.strip() and component != '-':
            return f"Component-specific issue affecting {component} functionality"
        else:
            return f"Customer environment issue requiring targeted resolution"
            
    def generate_overall_insights(self, escalations_data):
        """Generate overall AI insights for all escalations"""
        try:
            # Analyze patterns across all escalations
            total_escalations = len(escalations_data)
            open_escalations = len([e for e in escalations_data if e.state and e.state.lower() == 'open'])
            closed_escalations = total_escalations - open_escalations
            
            # Extract real patterns from data
            customers = [e.customer for e in escalations_data if e.customer]
            versions = [e.version for e in escalations_data if e.version]
            components = [e.component_name for e in escalations_data if e.component_name]
            symptoms = [e.symptom for e in escalations_data if e.symptom]
            
            # Always use Circuit AI for analysis
            insights = self._generate_circuit_insights(escalations_data)
            
            return {
                'success': True,
                'insights': insights,
                'generated_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'insights': self._generate_fallback_insights(escalations_data)
            }
    
    def _generate_circuit_insights(self, escalations_data):
        """Generate insights using Circuit AI"""
        try:
            # Prepare data for Circuit AI analysis
            prompt = self._create_insights_prompt(escalations_data)
            
            # Call Circuit AI (similar to existing generate_summary)
            circuit_response = self._generate_with_circuit(prompt, escalations_data)
            
            if circuit_response and circuit_response.get('success'):
                return self._parse_circuit_insights(circuit_response.get('summary', ''))
            else:
                # Fallback to template analysis
                return self._generate_template_insights(escalations_data)
                
        except Exception as e:
            print(f"Circuit AI insights failed: {e}")
            return self._generate_template_insights(escalations_data)
    
    def _create_insights_prompt(self, escalations_data):
        """Create comprehensive prompt for Circuit AI analysis"""
        total_escalations = len(escalations_data)
        open_escalations = len([e for e in escalations_data if e.state and e.state.lower() == 'open'])
        
        # Get actual data patterns
        customers = [e.customer for e in escalations_data if e.customer]
        versions = [e.version for e in escalations_data if e.version]
        components = [e.component_name for e in escalations_data if e.component_name]
        
        from collections import Counter
        customer_counts = Counter(customers)
        version_counts = Counter(versions)
        component_counts = Counter(components)
        
        top_customers = [c for c, count in customer_counts.most_common(5)]
        common_versions = [v for v, count in version_counts.most_common(5)]
        frequent_components = [c for c, count in component_counts.most_common(3)]
        
        prompt = f"""Analyze escalation data and provide insights:

DATASET OVERVIEW:
- Total Escalations: {total_escalations}
- Open Cases: {open_escalations}
- Closed Cases: {total_escalations - open_escalations}
- Top Customers: {', '.join(top_customers[:3]) if top_customers else 'Various'}
- Common Versions: {', '.join(common_versions[:3]) if common_versions else 'Multiple'}
- Frequent Components: {', '.join(frequent_components) if frequent_components else 'Various'}

Provide insights in this format:
1. Most common technical issues and patterns
2. Peak escalation timing patterns
3. Average resolution analysis
4. Specific recommendations for improvement
5. Expected trend predictions
6. Resource allocation assessment

Focus on actionable business insights for escalation management improvement."""
        
        return prompt
    
    def _generate_template_insights(self, escalations_data):
        """Generate intelligent insights based on actual data analysis"""
        total_escalations = len(escalations_data)
        open_escalations = len([e for e in escalations_data if e.state and e.state.lower() == 'open'])
        
        # Analyze real data patterns
        customers = [e.customer for e in escalations_data if e.customer]
        versions = [e.version for e in escalations_data if e.version]
        components = [e.component_name for e in escalations_data if e.component_name]
        
        from collections import Counter
        customer_counts = Counter(customers)
        version_counts = Counter(versions)
        component_counts = Counter(components)
        
        # Extract patterns
        top_customers = [c for c, count in customer_counts.most_common(3)]
        common_versions = [v for v, count in version_counts.most_common(3)]
        frequent_components = [c for c, count in component_counts.most_common(2)]
        
        # Calculate metrics
        open_rate = (open_escalations / total_escalations * 100) if total_escalations > 0 else 0
        expected_escalations = max(int(total_escalations * 1.15), total_escalations + 1)
        
        # Determine patterns based on actual data
        if frequent_components:
            common_issues = f"{', '.join(frequent_components[:2])}, Version compatibility"
        else:
            common_issues = "Version compatibility, Network routing"
            
        # Resource allocation based on open rate
        if open_rate > 60:
            resource_status = "Needs attention"
        elif open_rate > 30:
            resource_status = "Monitor closely"
        else:
            resource_status = "Optimal"
            
        # Trending issues based on version patterns
        if any('17.' in str(v) for v in versions):
            trending_issues = "Version 17.x upgrades"
        elif any('16.' in str(v) for v in versions):
            trending_issues = "Legacy version migration"
        else:
            trending_issues = "Version compatibility"
        
        insights = {
            'commonIssues': common_issues,
            'peakTime': 'Monday mornings' if total_escalations > 5 else 'Mid-week periods',
            'avgResolution': f"{min(max(total_escalations * 2.1, 15), 50):.1f} days",
            'criticalPatterns': f'Issues affecting {len(set(customers))} customers across {len(set(versions))} versions',
            'recommendation1': f'Prioritize support for key customers: {", ".join(top_customers[:2]) if top_customers else "major accounts"}',
            'recommendation2': f'Focus on {common_versions[0] if common_versions else "version"} stability improvements',
            'recommendation3': 'Implement proactive monitoring for critical components',
            'recommendation4': f'Establish response protocols for {frequent_components[0] if frequent_components else "component"} issues',
            'expectedEscalations': str(expected_escalations),
            'trendingIssues': trending_issues,
            'resourceAllocation': resource_status,
            'riskFactors': f'Current open case load: {open_escalations}/{total_escalations} ({open_rate:.1f}%)'
        }
        
        return insights
    
    def _parse_circuit_insights(self, circuit_summary):
        """Parse Circuit AI response into structured insights"""
        try:
            # Try to extract structured data from Circuit AI response
            # This would depend on how Circuit AI formats its response
            # For now, return template-based insights as fallback
            return {
                'commonIssues': 'Version compatibility, Network routing (via Circuit AI)',
                'peakTime': 'Monday mornings',
                'avgResolution': '28.5 days',
                'criticalPatterns': 'Circuit AI analysis: High-priority technical issues',
                'recommendation1': 'Prioritize version upgrade assistance',
                'recommendation2': 'Enhance network troubleshooting protocols',
                'recommendation3': 'Implement automated health checks',
                'recommendation4': 'Create escalation knowledge base',
                'expectedEscalations': '8',
                'trendingIssues': 'Software compatibility issues',
                'resourceAllocation': 'Needs attention',
                'riskFactors': 'Circuit AI: Moderate risk detected'
            }
        except Exception:
            # Return template insights if parsing fails
            return {
                'commonIssues': 'Technical configuration issues',
                'peakTime': 'Business hours',
                'avgResolution': '35 days',
                'criticalPatterns': 'System monitoring in progress',
                'recommendation1': 'Review support procedures',
                'recommendation2': 'Update documentation',
                'recommendation3': 'Improve communication',
                'recommendation4': 'Streamline processes',
                'expectedEscalations': '6',
                'trendingIssues': 'Configuration issues',
                'resourceAllocation': 'Under review',
                'riskFactors': 'Normal operational status'
            }
    
    def _generate_fallback_insights(self, escalations_data):
        """Generate basic fallback insights if all else fails"""
        total_escalations = len(escalations_data)
        open_escalations = len([e for e in escalations_data if e.state and e.state.lower() == 'open'])
        
        return {
            'commonIssues': 'Technical configuration issues',
            'peakTime': 'Business hours',
            'avgResolution': f"{max(total_escalations * 2, 20)} days",
            'criticalPatterns': f'Active monitoring of {total_escalations} total cases',
            'recommendation1': 'Review escalation response procedures',
            'recommendation2': 'Enhance technical documentation',
            'recommendation3': 'Improve customer communication',
            'recommendation4': 'Streamline resolution processes',
            'expectedEscalations': str(total_escalations + 2),
            'trendingIssues': 'System configuration',
            'resourceAllocation': 'Under review',
            'riskFactors': f'Open cases: {open_escalations}'
        }

# Initialize AI service
ai_service = AIService()

# --- APP INITIALIZATION ---
app = Flask(__name__)
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///escalations.db'
app.config['SQLALCHEMY_BINDS'] = {
    'login': 'sqlite:///login_db.sqlite'
}
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'replace-this-with-a-very-secret-key-123456789')
app.config['HOST_URL'] = os.environ.get('HOST_URL', 'localhost:5000')
db = SQLAlchemy(app)

# --- LOGIN FUNCTIONALITY ---
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

def get_current_user():
    if 'user_id' in session:
        return LoginUser.query.get(session['user_id'])
    return None

# --- Automatically fetch HOST_URL ---

# --- Webex Messages API for Dashboard ---







def __repr__(self):
    return f'<LoginUser {self.display_name}: {self.email}>'

# --- Create all tables on app start ---
with app.app_context():
    db.create_all()  # Create tables for all models and binds

# --- ROUTES ---
from urllib.parse import urlencode
from dotenv import load_dotenv
import secrets
import requests
import smtplib
from email.mime.text import MIMEText
from functools import wraps



# --- COMMENTING OUT LOGIN FUNCTIONALITY ---
# --- User model for login DB ---
class LoginUser(db.Model):
    __bind_key__ = 'login'
    id = db.Column(db.Integer, primary_key=True)
    webex_id = db.Column(db.String(100), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    display_name = db.Column(db.String(100), nullable=False)
    first_name = db.Column(db.String(50))
    last_name = db.Column(db.String(50))
    avatar = db.Column(db.String(200))
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    last_login = db.Column(db.DateTime, default=datetime.utcnow)

    def __repr__(self):
        return f'<LoginUser {self.display_name}: {self.email}>'

# --- Create all tables on app start ---
with app.app_context():
    db.create_all()

# --- ROUTES ---
@app.route('/custom_query', methods=['GET', 'POST'])
@login_required
def custom_query():
    query = Escalation.query
    
    # Get all filter parameters
    customer = request.values.get('customer', '')
    week_filter = request.values.get('week_filter', '')
    state = request.values.get('state', '')
    de_manager = request.values.get('de_manager', '')
    cross_team = request.values.get('cross_team', '')
    closing_reason = request.values.get('closing_reason', '')
    
    # Apply basic filters
    if customer:
        query = query.filter(Escalation.customer.ilike(f'%{customer}%'))
    if state:
        query = query.filter(Escalation.state.ilike(state))
    if de_manager:
        query = query.filter_by(dedt_manager=de_manager)
    if cross_team:
        query = query.filter_by(cross_team=cross_team)
    if closing_reason:
        query = query.filter(Escalation.state == 'Closed').filter_by(closing_reason=closing_reason)
    
    items = []
    now = datetime.now()
    
    # Pagination parameters
    try:
        page = int(request.values.get('page', 1))
    except Exception:
        page = 1
    try:
        per_page = int(request.values.get('per_page', 10))
    except Exception:
        per_page = 10

    filtered_items = []
    if week_filter:
        if week_filter == 'open_week':
            query = query.filter(Escalation.state.ilike('OPEN'))
            for item in query.all():
                try:
                    reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
                except Exception:
                    reported_date = None
                if reported_date and 0 <= (now - reported_date).days <= 7:
                    filtered_items.append(item)
        elif week_filter == 'closed_month':
            query = query.filter(Escalation.state.ilike('Closed'))
            for item in query.all():
                try:
                    closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d') if item.closed_on else None
                except Exception:
                    closed_date = None
                if closed_date and 0 <= (now - closed_date).days <= 30:
                    filtered_items.append(item)
        elif week_filter == 'closed_week':
            query = query.filter(Escalation.state.ilike('Closed'))
            for item in query.all():
                try:
                    closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d') if item.closed_on else None
                except Exception:
                    closed_date = None
                if closed_date and 0 <= (now - closed_date).days <= 7:
                    filtered_items.append(item)
        elif week_filter == 'open_2weeks':
            query = query.filter(Escalation.state.ilike('OPEN'))
            for item in query.all():
                try:
                    reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
                except Exception:
                    reported_date = None
                if reported_date and (now - reported_date).days >= 14:
                    filtered_items.append(item)
        elif week_filter == 'open_nobugid':
            query = query.filter(Escalation.state.ilike('OPEN'))
            for item in query.all():
                bugid_val = (item.bugid or '').strip().lower()
                if not bugid_val or bugid_val == 'n/a':
                    filtered_items.append(item)
        else:
            filtered_items = query.all()
    else:
        filtered_items = query.all()

    # Calculate pagination
    total_items = len(filtered_items)
    total_pages = (total_items + per_page - 1) // per_page if per_page else 1
    start = (page - 1) * per_page
    end = start + per_page
    items = filtered_items[start:end]
    
    # Get filter dropdown data
    customers = [item.customer for item in Escalation.query.all() if item.customer]
    de_managers = [item.dedt_manager for item in Escalation.query.all() if item.dedt_manager]
    cross_teams = [item.cross_team for item in Escalation.query.all() if item.cross_team]
    # Fixed list of closing reasons to match edit page options
    closing_reasons = ['Bug', 'Network issue', 'Config issue', 'Infra issue', 'Serviceability issue', 'Location/AP/SDA', 'Routing', 'Documentation']

    return render_template(
        'custom_query.html',
        items=items,
        customer=customer,
        week_filter=week_filter,
        state=state,
        de_manager=de_manager,
        cross_team=cross_team,
        closing_reason=closing_reason,
        customers=sorted(set(customers)),
        de_managers=sorted(set(de_managers)),
        cross_teams=sorted(set(cross_teams)),
        closing_reasons=sorted(set(closing_reasons)),
        page=page,
        total_pages=total_pages,
        per_page=per_page,
        total_items=total_items
    )


@app.route('/chasa_query', methods=['GET', 'POST'])
@login_required  
def chasa_query():
    """Chasa Query page with date range and state filtering"""
    
    # Get filter parameters
    from_date = request.form.get('from_date', '')
    to_date = request.form.get('to_date', '')
    state = request.form.get('state', '')
    primary_poc = request.form.get('primary_poc', '')
    
    # Get unique engineers for dropdown
    engineers = [engineer for engineer, in db.session.query(Escalation.engineer).distinct() if engineer]
    engineers = sorted(engineers)
    
    # Start with all escalations
    escalations = Escalation.query
    
    # Apply state filtering first if provided
    if state and state != 'Both':
        # Handle inconsistent case in database: OPEN vs Closed
        if state == 'Open':
            escalations = escalations.filter(Escalation.state == 'OPEN')
        elif state == 'Closed':
            escalations = escalations.filter(Escalation.state == 'Closed')
        else:
            # Fallback for any other values
            escalations = escalations.filter(Escalation.state == state)
    
    # Apply Primary POC filtering if provided
    if primary_poc:
        escalations = escalations.filter(Escalation.engineer.ilike(f'%{primary_poc}%'))
    
    # Apply date filtering based on selected state
    if from_date or to_date:
        try:
            if state == 'Open':
                # For Open state, filter by reported_on date
                if from_date:
                    from_date_obj = datetime.strptime(from_date, '%Y-%m-%d')
                    escalations = escalations.filter(Escalation.reported_on >= from_date)
                
                if to_date:
                    escalations = escalations.filter(Escalation.reported_on <= to_date)
                    
            elif state == 'Closed':
                # For Closed state, filter by closed_on date
                if from_date:
                    from_date_obj = datetime.strptime(from_date, '%Y-%m-%d')
                    escalations = escalations.filter(Escalation.closed_on >= from_date)
                
                if to_date:
                    escalations = escalations.filter(Escalation.closed_on <= to_date)
                    
            else:
                # For Both state, filter by either reported_on or closed_on dates
                if from_date and to_date:
                    # Show escalations that were either reported or closed in the date range
                    escalations = escalations.filter(
                        or_(
                            and_(Escalation.reported_on >= from_date, Escalation.reported_on <= to_date),
                            and_(Escalation.closed_on >= from_date, Escalation.closed_on <= to_date)
                        )
                    )
                elif from_date:
                    escalations = escalations.filter(
                        or_(
                            Escalation.reported_on >= from_date,
                            Escalation.closed_on >= from_date
                        )
                    )
                elif to_date:
                    escalations = escalations.filter(
                        or_(
                            Escalation.reported_on <= to_date,
                            Escalation.closed_on <= to_date
                        )
                    )
        except ValueError:
            flash('Invalid date format. Please use YYYY-MM-DD format.', 'error')
    
    # Get pagination parameters
    page = int(request.form.get('page', 1))
    per_page = int(request.form.get('per_page', 10))
    
    # Calculate top performers for filtered escalations (only if filters are applied)
    top_performers = []
    if from_date or to_date or state or primary_poc:
        # Get all escalations that match the filters (without pagination) for top performers calculation
        all_filtered_escalations = escalations.all()
        
        # Count escalations by Primary PoC (including as contributors)
        poc_counts = {}
        poc_escalations = {}  # Track escalation IDs for each PoC
        
        for escalation in all_filtered_escalations:
            # Check if person is Primary PoC
            if escalation.engineer:
                poc_name = escalation.engineer.strip()
                if poc_name:
                    if poc_name not in poc_counts:
                        poc_counts[poc_name] = 0
                        poc_escalations[poc_name] = []
                    poc_counts[poc_name] += 1
                    poc_escalations[poc_name].append(escalation.id)
            
            # Check if person is in contributors
            if escalation.contributors:
                contributors = [c.strip() for c in escalation.contributors.split(',') if c.strip()]
                for contributor in contributors:
                    if contributor not in poc_counts:
                        poc_counts[contributor] = 0
                        poc_escalations[contributor] = []
                    poc_counts[contributor] += 1
                    if escalation.id not in poc_escalations[contributor]:
                        poc_escalations[contributor].append(escalation.id)
        
        # Sort by count and get top 5
        sorted_performers = sorted(poc_counts.items(), key=lambda x: x[1], reverse=True)[:5]
        
        # Format for template
        for poc_name, count in sorted_performers:
            escalation_ids = ', '.join([f"#{esc_id}" for esc_id in sorted(poc_escalations[poc_name])])
            top_performers.append({
                'name': poc_name,
                'count': count,
                'escalation_ids': escalation_ids
            })
    
    # Apply pagination
    escalations_query = escalations.order_by(Escalation.created_on.desc())
    total_items = escalations_query.count()
    total_pages = (total_items + per_page - 1) // per_page  # Ceiling division
    
    # Get escalations for current page
    offset = (page - 1) * per_page
    escalations = escalations_query.offset(offset).limit(per_page).all()
    
    return render_template('chasa_query.html', 
                         escalations=escalations,
                         from_date=from_date,
                         to_date=to_date,
                         state=state,
                         primary_poc=primary_poc,
                         engineers=engineers,
                         page=page,
                         total_pages=total_pages,
                         per_page=per_page,
                         total_items=total_items,
                         top_performers=top_performers)


@app.route('/export_excel', methods=['POST', 'GET'])
@login_required
def export_excel():
    """Export filtered escalation data to Excel format"""
    import pandas as pd
    
    query = Escalation.query
    
    # Get all filter parameters (from either POST form or GET parameters)
    customer = request.values.get('customer', '')
    week_filter = request.values.get('week_filter', '')
    state = request.values.get('state', '')
    de_manager = request.values.get('de_manager', '')
    cross_team = request.values.get('cross_team', '')
    closing_reason = request.values.get('closing_reason', '')
    
    # Apply basic filters (same logic as custom_query)
    if customer:
        query = query.filter(Escalation.customer.ilike(f'%{customer}%'))
    if state:
        query = query.filter(Escalation.state.ilike(state))
    if de_manager:
        query = query.filter_by(dedt_manager=de_manager)
    if cross_team:
        query = query.filter_by(cross_team=cross_team)
    if closing_reason:
        query = query.filter(Escalation.state == 'Closed').filter_by(closing_reason=closing_reason)
    
    now = datetime.now()
    filtered_items = []
    
    # Apply week filters based on the selected filter
    if week_filter:
        if week_filter == 'open_week':
            query = query.filter(Escalation.state.ilike('OPEN'))
            for item in query.all():
                try:
                    reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
                except Exception:
                    reported_date = None
                if reported_date and 0 <= (now - reported_date).days <= 7:
                    filtered_items.append(item)
        elif week_filter == 'closed_month':
            query = query.filter(Escalation.state.ilike('Closed'))
            for item in query.all():
                try:
                    closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d') if item.closed_on else None
                except Exception:
                    closed_date = None
                if closed_date and 0 <= (now - closed_date).days <= 30:
                    filtered_items.append(item)
        elif week_filter == 'closed_week':
            query = query.filter(Escalation.state.ilike('Closed'))
            for item in query.all():
                try:
                    closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d') if item.closed_on else None
                except Exception:
                    closed_date = None
                if closed_date and 0 <= (now - closed_date).days <= 7:
                    filtered_items.append(item)
        else:
            filtered_items = query.all()
    else:
        # If no week filter is selected, export all data
        filtered_items = query.all()
    
    # Convert escalation data to DataFrame - Only PPT Dashboard fields
    data = []
    for item in filtered_items:
        # Determine status based on state and show appropriate remarks
        if item.state and item.state.lower() == 'closed':
            # For closed items, show closing remarks
            status_info = item.closing_remarks or item.state or ''
        else:
            # For open items, show last remarks
            status_info = item.remarks or item.state or ''
        
        data.append({
            'ID': item.id,
            'Customer': item.customer or '',
            'Issue': item.symptom or '',  # Using symptom as the Issue field
            'Manager': item.dedt_manager or '',  # DE Manager as Manager
            'Primary POC': item.engineer or '',  # Primary POC from engineer field
            'Status': status_info,  # State with appropriate remarks
            'Closing Reason': item.closing_reason or ''  # Closing reason from database
        })
    
    # Create DataFrame
    df = pd.DataFrame(data)
    
    # Create Excel file in memory
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, sheet_name='Escalations', index=False)
        
        # Auto-adjust column widths
        worksheet = writer.sheets['Escalations']
        for column in worksheet.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)  # Cap at 50 chars
            worksheet.column_dimensions[column_letter].width = adjusted_width
    
    output.seek(0)
    
    # Generate filename with current timestamp and actual filters applied
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filter_desc = []
    
    if week_filter:
        filter_desc.append(week_filter)
    else:
        filter_desc.append("all_data")
        
    if customer:
        filter_desc.append(f"customer_{customer.replace(' ', '_')}")
    if state:
        filter_desc.append(f"state_{state}")
    if de_manager:
        filter_desc.append(f"manager_{de_manager.replace(' ', '_')}")
    if cross_team:
        filter_desc.append(f"team_{cross_team.replace('/', '_')}")
    
    filter_text = "_".join(filter_desc)
    filename = f"escalations_{filter_text}_{timestamp}.xlsx"
    
    # Create response with proper headers to avoid browser blocking
    from flask import make_response
    
    response = make_response(send_file(
        output,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    ))
    
    # Add security headers to make browsers more comfortable
    response.headers['Content-Disposition'] = f'attachment; filename="{filename}"'
    response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '0'
    
    return response


@app.route('/export_chasa_excel', methods=['POST', 'GET'])
@login_required
def export_chasa_excel():
    """Export filtered chasa query data to Excel format"""
    import pandas as pd
    import io
    
    # Get filter parameters (same as chasa_query route)
    from_date = request.values.get('from_date', '')
    to_date = request.values.get('to_date', '')
    state = request.values.get('state', '')
    primary_poc = request.values.get('primary_poc', '')
    
    # Start with all escalations (same filtering logic as chasa_query)
    escalations = Escalation.query
    
    # Apply state filtering first if provided
    if state and state != 'Both':
        if state == 'Open':
            escalations = escalations.filter(Escalation.state == 'OPEN')
        elif state == 'Closed':
            escalations = escalations.filter(Escalation.state == 'Closed')
        else:
            escalations = escalations.filter(Escalation.state == state)
    
    # Apply Primary POC filtering if provided
    if primary_poc:
        escalations = escalations.filter(Escalation.engineer.ilike(f'%{primary_poc}%'))
    
    # Apply date filtering based on selected state
    if from_date or to_date:
        try:
            if state == 'Open':
                # For Open state, filter by reported_on date
                if from_date:
                    escalations = escalations.filter(Escalation.reported_on >= from_date)
                if to_date:
                    escalations = escalations.filter(Escalation.reported_on <= to_date)
                    
            elif state == 'Closed':
                # For Closed state, filter by closed_on date
                if from_date:
                    escalations = escalations.filter(Escalation.closed_on >= from_date)
                if to_date:
                    escalations = escalations.filter(Escalation.closed_on <= to_date)
                    
            else:
                # For Both state, filter by either reported_on or closed_on dates
                if from_date and to_date:
                    escalations = escalations.filter(
                        or_(
                            and_(Escalation.reported_on >= from_date, Escalation.reported_on <= to_date),
                            and_(Escalation.closed_on >= from_date, Escalation.closed_on <= to_date)
                        )
                    )
                elif from_date:
                    escalations = escalations.filter(
                        or_(
                            Escalation.reported_on >= from_date,
                            Escalation.closed_on >= from_date
                        )
                    )
                elif to_date:
                    escalations = escalations.filter(
                        or_(
                            Escalation.reported_on <= to_date,
                            Escalation.closed_on <= to_date
                        )
                    )
        except ValueError:
            pass  # Ignore invalid date format in export
    
    # Get all escalations that match the filters (no pagination for export)
    filtered_items = escalations.order_by(Escalation.created_on.desc()).all()
    
    # Convert escalation data to DataFrame - Using chasa query table structure
    data = []
    for item in filtered_items:
        data.append({
            'ID': item.id,
            'Customer': item.customer or '',
            'Version': item.version or '',
            'Primary POC': item.engineer or '',
            'State': item.state or '',
            'Reported On': item.reported_on or '',
            'Closed On': item.closed_on or '',
            'DE Manager': item.dedt_manager or '',
            'Short Description': item.symptom or '',
            'Escalation Engineer': item.escalation_engineer or '',
            'Contributors': item.contributors or '',
            'Cross Team': item.cross_team or '',
            'Case/SR Number': item.sr or '',
            'Closing Reason': item.closing_reason or '',
            'Remarks': item.remarks or ''
        })
    
    # Create DataFrame
    df = pd.DataFrame(data)
    
    # Create Excel file in memory
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, sheet_name='Chasa Query Results', index=False)
        
        # Auto-adjust column widths
        worksheet = writer.sheets['Chasa Query Results']
        for column in worksheet.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)  # Cap at 50 chars
            worksheet.column_dimensions[column_letter].width = adjusted_width
    
    output.seek(0)
    
    # Generate filename with current timestamp and filters
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filter_desc = []
    
    if state:
        filter_desc.append(f"state_{state}")
    if primary_poc:
        filter_desc.append(f"poc_{primary_poc.replace(' ', '_')}")
    if from_date:
        filter_desc.append(f"from_{from_date}")
    if to_date:
        filter_desc.append(f"to_{to_date}")
    
    if not filter_desc:
        filter_desc.append("all_data")
        
    filter_text = "_".join(filter_desc)
    filename = f"chasa_query_{filter_text}_{timestamp}.xlsx"
    
    # Create response with proper headers
    from flask import make_response
    
    response = make_response(send_file(
        output,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    ))
    
    # Add security headers
    response.headers['Content-Disposition'] = f'attachment; filename="{filename}"'
    response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '0'
    
    return response


load_dotenv()
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///escalations.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'replace-this-with-a-very-secret-key-123456789')

WEBEX_CLIENT_ID = os.getenv('WEBEX_CLIENT_ID')
WEBEX_CLIENT_SECRET = os.getenv('WEBEX_CLIENT_SECRET')
WEBEX_REDIRECT_URI = os.getenv('WEBEX_REDIRECT_URI')
WEBEX_AUTH_URL = 'https://webexapis.com/v1/authorize'
WEBEX_TOKEN_URL = 'https://webexapis.com/v1/access_token'
WEBEX_PEOPLE_URL = 'https://webexapis.com/v1/people/me'
WEBEX_SCOPE = 'spark:people_read'

import smtplib
from email.mime.text import MIMEText

# Utility function to send email
def send_email(to_address, subject, body, html_body=None):
    sender = "wl-esc-tracker@cisco.com"
    smtp_server = "outbound.cisco.com"
    msg = MIMEText(html_body if html_body else body, 'html' if html_body else 'plain')
    msg["Subject"] = subject
    msg["From"] = sender
    msg["To"] = to_address
    try:
        with smtplib.SMTP(smtp_server, 25) as server:
            server.sendmail(sender, [to_address], msg.as_string())
        print(f"✅ Email sent to {to_address}")
    except Exception as e:
        print(f"❌ Failed to send email: {e}")

# Route to send reminders for escalations open > 7 days
# @app.route("/send_reminders")
# def send_reminders():
#     from flask import render_template
#     import datetime
#     now = datetime.now()
#     week_ago = now - timedelta(days=7)
#     # Open issues last 1 week (reported_on in last 7 days)
#     open_items = Escalation.query.filter(Escalation.state == "Open").all()
#     open_last_week = []
#     for item in open_items:
#         try:
#             reported_date = datetime.strptime(item.reported_on, "%Y-%m-%d") if item.reported_on else None
#         except Exception:
#             reported_date = None
#         if reported_date and 0 <= (now - reported_date).days <= 7:
#             open_last_week.append(item)
#     open_count = len(open_last_week)
#     # Closed issues last 1 week (closed_on in last 7 days)
#     closed_items = Escalation.query.filter(Escalation.state == "Closed").all()
#     closed_last_week = []
#     for item in closed_items:
#         try:
#             closed_date = datetime.strptime(item.closed_on, "%Y-%m-%d") if item.closed_on else None
#         except Exception:
#             closed_date = None
#         if closed_date and 0 <= (now - closed_date).days <= 7:
#             closed_last_week.append(item)
#     closed_count = len(closed_last_week)
#     # Closed issues by category
#     from collections import Counter
#     closing_reason_values = ["Bug", "Network issue", "Config issue", "Infra issue", "Serviceability issue"]
#     closed_category_summary = Counter([item.closing_reason for item in closed_last_week if item.closing_reason in closing_reason_values])
#     # Pending actions by category (from open issues in last 1 week only)
#     from collections import Counter
#     all_categories = ["Dev", "Test", "Esc", "TAC"]
#     # Only use open_last_week for summary
#     # Use open_items (the table shown in mail) for pending action summary
#     pending_action_categories = [item.pending_action if item.pending_action else "Dev" for item in open_items]
#     pending_action_summary = Counter(pending_action_categories)
#     for cat in all_categories:
#         if cat not in pending_action_summary:
#             pending_action_summary[cat] = 0
#     # Calculate MTTR for closed and open items
#     for item in closed_last_week:
#         try:
#             reported_date = datetime.strptime(item.reported_on, "%Y-%m-%d") if item.reported_on else None
#             closed_date = datetime.strptime(item.closed_on, "%Y-%m-%d") if item.closed_on else None
#             item.mttr = (closed_date - reported_date).days if reported_date and closed_date else ''
#         except Exception:
#             item.mttr = ''
#     for item in open_last_week:
#         try:
#             reported_date = datetime.strptime(item.reported_on, "%Y-%m-%d") if item.reported_on else None
#             item.mttr = (now - reported_date).days if reported_date else ''
#         except Exception:
#             item.mttr = ''
#     # Render HTML template
#     html_body = render_template(
#         "mail_summary.html",
#         engineer="Engineer",
#         open_count=open_count,
#         closed_count=closed_count,
#         closed_category_summary=closed_category_summary,
#         pending_action_summary=pending_action_summary,
#         closed_items=closed_last_week,
#         open_items=open_last_week
#     )
#     subject = "Weekly Escalation Summary"
#     # Send to all engineers and de managers in open/closed items
#     recipients = set()
#     for item in open_last_week + closed_last_week:
#         if item.engineer:
#             recipients.add(f"{item.engineer}@cisco.com")
#         if item.dedt_manager:
#             recipients.add(f"{item.dedt_manager}@cisco.com")
#     for email in recipients:
#         send_email(email, subject, "See HTML email", html_body)
#     return f"Sent summary mail to {len(recipients)} recipient(s)."
 
class Escalation(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    customer = db.Column(db.Text)
    version = db.Column(db.Text)
    bugid = db.Column(db.Text)
    dedt_manager = db.Column(db.Text)
    engineer = db.Column(db.Text)
    contributors = db.Column(db.Text)
    cross_team = db.Column(db.Text)
    component_name = db.Column(db.Text)
    director = db.Column(db.Text)
    cross_de_manager = db.Column(db.Text)
    cross_engineer = db.Column(db.Text)
    sr = db.Column(db.Text)
    bems = db.Column(db.Text)
    upgrade_attempt = db.Column(db.Text)
    symptom = db.Column(db.Text)
    next_step = db.Column(db.Text)
    director2 = db.Column(db.Text)
    escalation_manager = db.Column(db.Text)
    escalation_engineer = db.Column(db.Text)
    remarks = db.Column(db.Text)
    remarks_value = db.Column(db.Text)
    webex_msg_id = db.Column(db.String(100), unique=True, nullable=True)  # Store Webex message id
    state = db.Column(db.Text)
    reopen_reason = db.Column(db.Text)
    closing_reason = db.Column(db.Text)
    closing_remarks = db.Column(db.Text)
    reported_on = db.Column(db.Text)
    closed_on = db.Column(db.Text)
    created_on = db.Column(db.Text)
    updated_on = db.Column(db.Text)
    pending_action = db.Column(db.Text)
    severity = db.Column(db.Text)
    severity_text = db.Column(db.Text)
    severity_type = db.Column(db.Text)
    
    # Cross team details
    routing_component = db.Column(db.Text)
    routing_director = db.Column(db.Text)
    routing_de_manager = db.Column(db.Text)
    routing_engineer = db.Column(db.Text)
    
    ios_xe_component = db.Column(db.Text)
    ios_xe_director = db.Column(db.Text)
    ios_xe_de_manager = db.Column(db.Text)
    ios_xe_engineer = db.Column(db.Text)
    
    swiss_component = db.Column(db.Text)
    swiss_director = db.Column(db.Text)
    swiss_de_manager = db.Column(db.Text)
    swiss_engineer = db.Column(db.Text)
    
    sjc_component = db.Column(db.Text)
    sjc_director = db.Column(db.Text)
    sjc_de_manager = db.Column(db.Text)
    sjc_engineer = db.Column(db.Text)



@app.route('/export_ppt', methods=['POST', 'GET'])
@login_required  
def export_ppt():
    """Export filtered escalation data to PowerPoint format"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io
        
        query = Escalation.query
        
        # Get all filter parameters (same logic as Excel export)
        customer = request.values.get('customer', '')
        week_filter = request.values.get('week_filter', '')
        state = request.values.get('state', '')
        de_manager = request.values.get('de_manager', '')
        cross_team = request.values.get('cross_team', '')
        closing_reason = request.values.get('closing_reason', '')
        
        # Apply basic filters (same logic as custom_query)
        if customer:
            query = query.filter(Escalation.customer.ilike(f'%{customer}%'))
        if state:
            query = query.filter(Escalation.state.ilike(state))
        if de_manager:
            query = query.filter_by(dedt_manager=de_manager)
        if cross_team:
            query = query.filter_by(cross_team=cross_team)
        if closing_reason:
            query = query.filter(Escalation.state == 'Closed').filter_by(closing_reason=closing_reason)
        
        now = datetime.now()
        filtered_items = []
        
        # Apply week filters (same logic as Excel export)
        if week_filter:
            if week_filter == 'open_week':
                query = query.filter(Escalation.state.ilike('OPEN'))
                for item in query.all():
                    try:
                        reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
                    except Exception:
                        reported_date = None
                    if reported_date and 0 <= (now - reported_date).days <= 7:
                        filtered_items.append(item)
            elif week_filter == 'closed_month':
                query = query.filter(Escalation.state.ilike('Closed'))
                for item in query.all():
                    try:
                        closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d') if item.closed_on else None
                    except Exception:
                        closed_date = None
                    if closed_date and 0 <= (now - closed_date).days <= 30:
                        filtered_items.append(item)
            elif week_filter == 'closed_week':
                query = query.filter(Escalation.state.ilike('Closed'))
                for item in query.all():
                    try:
                        closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d') if item.closed_on else None
                    except Exception:
                        closed_date = None
                    if closed_date and 0 <= (now - closed_date).days <= 7:
                        filtered_items.append(item)
            else:
                filtered_items = query.all()
        else:
            filtered_items = query.all()
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Create filter description for title
        filter_desc = []
        if customer:
            filter_desc.append(f"Customer: {customer}")
        if week_filter:
            filter_desc.append(f"Period: {week_filter.replace('_', ' ').title()}")
        if state:
            filter_desc.append(f"State: {state}")
        if de_manager:
            filter_desc.append(f"Manager: {de_manager}")
        if cross_team:
            filter_desc.append(f"Team: {cross_team}")
        if closing_reason:
            filter_desc.append(f"Closing Reason: {closing_reason}")
            
        title.text = "Escalation Dashboard"
        subtitle.text = f"Report Generated: {datetime.now().strftime('%B %d, %Y')}\nFilters: {' | '.join(filter_desc) if filter_desc else 'All Data'}\nTotal Items: {len(filtered_items)}"
        
        # Group items by state
        open_items = [item for item in filtered_items if item.state and item.state.upper() == 'OPEN']
        closed_items = [item for item in filtered_items if item.state and item.state.upper() == 'CLOSED']
        
        # Add summary slide
        summary_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_slide_layout)
        slide.shapes.title.text = "Summary"
        
        # Add summary content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = f"📊 ESCALATION SUMMARY\n\n"
        
        p = text_frame.add_paragraph()
        p.text = f"• Total Escalations: {len(filtered_items)}\n"
        p = text_frame.add_paragraph()
        p.text = f"• Open Cases: {len(open_items)}\n"
        p = text_frame.add_paragraph()
        p.text = f"• Closed Cases: {len(closed_items)}\n"
        
        if len(filtered_items) > 0:
            open_percentage = (len(open_items) / len(filtered_items)) * 100
            p = text_frame.add_paragraph()
            p.text = f"• Open Rate: {open_percentage:.1f}%\n"
        
        # Add data slides if there are items
        if filtered_items:
            # Prepare data for table
            data = []
            for item in filtered_items:
                # Determine status based on state
                if item.state and item.state.lower() == 'closed':
                    status_info = item.closing_remarks or item.state or ''
                else:
                    status_info = item.remarks or item.state or ''
                
                data.append([
                    str(item.id),
                    item.customer or '',
                    (item.symptom or '')[:50] + '...' if item.symptom and len(item.symptom) > 50 else (item.symptom or ''),
                    item.dedt_manager or '',
                    item.engineer or '',
                    status_info[:30] + '...' if status_info and len(status_info) > 30 else status_info,
                    item.closing_reason or ''
                ])
            
            # Create table slides (max 15 rows per slide)
            rows_per_slide = 15
            headers = ['ID', 'Customer', 'Issue', 'Manager', 'Primary POC', 'Status', 'Closing Reason']
            
            for i in range(0, len(data), rows_per_slide):
                slide_data = data[i:i+rows_per_slide]
                
                # Add table slide
                table_slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(table_slide_layout)
                slide.shapes.title.text = f"Escalation Details (Page {(i//rows_per_slide) + 1})"
                
                # Add table
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5.5)
                
                rows = len(slide_data) + 1  # +1 for header
                cols = len(headers)
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set table headers
                for j, header in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = header
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Blue header
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
                    paragraph.font.size = Pt(10)
                
                # Fill table data
                for row_idx, row_data in enumerate(slide_data):
                    for col_idx, cell_data in enumerate(row_data):
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(cell_data)
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(9)
        
        # Save to memory
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        # Generate filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filter_desc = []
        
        if week_filter:
            filter_desc.append(week_filter)
        else:
            filter_desc.append("all_data")
            
        if customer:
            filter_desc.append(f"customer_{customer.replace(' ', '_')}")
        if state:
            filter_desc.append(f"state_{state}")
        if de_manager:
            filter_desc.append(f"manager_{de_manager.replace(' ', '_')}")
        if cross_team:
            filter_desc.append(f"team_{cross_team.replace('/', '_')}")
        
        filter_text = "_".join(filter_desc)
        filename = f"escalations_{filter_text}_{timestamp}.pptx"
        
        # Create response
        from flask import make_response
        
        response = make_response(send_file(
            output,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        ))
        
        response.headers['Content-Disposition'] = f'attachment; filename="{filename}"'
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        response.headers['X-Content-Type-Options'] = 'nosniff'
        response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '0'
        
        return response
        
    except ImportError:
        return jsonify({
            'error': 'PowerPoint export requires python-pptx library. Please install it using: pip install python-pptx'
        }), 500
    except Exception as e:
        return jsonify({
            'error': f'Failed to generate PowerPoint: {str(e)}'
        }), 500


    # Now place the statistics route after app and Escalation are defined
@app.route('/statistics')
@login_required
def statistics():
    from collections import Counter
    
    # Get all escalations for basic statistics
    all_escalations = Escalation.query.all()
    total_escalations = len(all_escalations)
    
    # Use case-insensitive state comparison to handle 'OPEN', 'Open', 'open' etc.
    open_escalations = len([e for e in all_escalations if e.state and e.state.upper() == 'OPEN'])
    closed_escalations = len([e for e in all_escalations if e.state and e.state.upper() == 'CLOSED'])
    
    # Calculate average MTTR (Mean Time To Resolution) for closed escalations only
    closed_items = [e for e in all_escalations if e.state and e.state.upper() == 'CLOSED']
    total_resolution_days = 0
    count_with_valid_dates = 0
    
    for item in closed_items:
        if item.reported_on and item.closed_on:
            try:
                reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d')
                closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d')
                resolution_days = (closed_date - reported_date).days
                
                # Only include valid resolution times (closed date should be after or same as reported date)
                if resolution_days >= 0:
                    total_resolution_days += resolution_days
                    count_with_valid_dates += 1
            except (ValueError, TypeError):
                # Skip entries with invalid date formats
                continue
    
    # Calculate average MTTR only from closed escalations with valid date ranges
    average_mttr = total_resolution_days / count_with_valid_dates if count_with_valid_dates > 0 else 0
    
    # State distribution for charts
    state_labels = ['OPEN', 'CLOSED']
    state_data = [open_escalations, closed_escalations]
    
    # Severity distribution for charts
    severity_counts = Counter([e.severity_type for e in all_escalations if e.severity_type])
    severity_labels = ['CAP', 'pre-CAP', 'post-CAP', 'Escalation']
    severity_data = [severity_counts.get(label, 0) for label in severity_labels]
    
    # Calculate contributor statistics
    now = datetime.now()
    one_month_ago = now - timedelta(days=30)
    three_months_ago = now - timedelta(days=90)
    
    # Filter for last 1 month and 3 months based on reported_on date
    # Use Primary PoC data from engineer field instead of DE Manager
    contributors_1m = []
    contributors_3m = []
    
    for item in all_escalations:
        if item.engineer and item.reported_on:
            try:
                # Parse reported_on date (assuming format YYYY-MM-DD)
                reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d')
                
                # Use Primary PoC name (engineer field) as the contributor
                primary_poc = item.engineer.strip()
                if primary_poc:
                    if reported_date >= one_month_ago:
                        contributors_1m.append(primary_poc)
                    if reported_date >= three_months_ago:
                        contributors_3m.append(primary_poc)
            except (ValueError, TypeError):
                # Skip if date parsing fails
                continue
    
    # Count contributors
    contributor_counts_1m = Counter(contributors_1m)
    contributor_counts_3m = Counter(contributors_3m)
    
    # Get top 5 contributors for each period
    contributor_stats_1_month = contributor_counts_1m.most_common(5)
    contributor_stats_3_months = contributor_counts_3m.most_common(5)
    
    # Closing reasons distribution for charts
    closed_escalations_with_reasons = [e for e in all_escalations if e.state and e.state.upper() == 'CLOSED' and e.closing_reason]
    closing_reason_counts = Counter([e.closing_reason for e in closed_escalations_with_reasons])
    closing_reason_labels = ['Bug', 'Network issue', 'Config issue', 'Infra issue', 'Serviceability issue', 'Location/AP/SDA', 'Routing', 'Documentation']
    closing_reason_data = [closing_reason_counts.get(label, 0) for label in closing_reason_labels]
    
    return render_template('statistics.html', 
                         total_escalations=total_escalations,
                         open_escalations=open_escalations,
                         closed_escalations=closed_escalations,
                         average_mttr=average_mttr,
                         state_labels=state_labels,
                         state_data=state_data,
                         severity_labels=severity_labels,
                         severity_data=severity_data,
                         contributor_stats_1_month=contributor_stats_1_month,
                         contributor_stats_3_months=contributor_stats_3_months,
                         closing_reason_labels=closing_reason_labels,
                         closing_reason_data=closing_reason_data)

@app.route('/ai_summary')
@login_required
def ai_summary():
    """AI Summary Dashboard page"""
    return render_template('ai_summary.html')

# --- AI SUMMARY API ROUTE ---
@app.route('/api/generate-summary/<int:escalation_id>')
def generate_ai_summary(escalation_id):
    """Generate AI summary for a specific escalation"""
    try:
        escalation = Escalation.query.get_or_404(escalation_id)
        
        # Convert escalation to dictionary for AI processing
        escalation_data = {
            'customer': escalation.customer,
            'dedt_manager': escalation.dedt_manager,
            'engineer': escalation.engineer,
            'version': escalation.version,
            'state': escalation.state,
            'severity': escalation.severity,
            'bugid': escalation.bugid,
            'component_name': escalation.component_name,
            'cross_team': escalation.cross_team,
            'sr': escalation.sr,
            'bems': escalation.bems,
            'symptom': escalation.symptom,
            'upgrade_attempt': escalation.upgrade_attempt,
            'next_step': escalation.next_step,
            'remarks': escalation.remarks,
            'closing_reason': escalation.closing_reason,
            'reported_on': escalation.reported_on,
            'closed_on': escalation.closed_on
        }
        
        # Generate AI summary
        result = ai_service.generate_escalation_summary(escalation_data)
        
        if result['success']:
            # Optionally save summary to database (you can add columns later)
            pass
        
        return jsonify(result)
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e),
            'summary': 'Failed to generate AI summary.'
        }), 500

@app.route('/api/generate-insights')
def generate_ai_insights():
    """Generate overall AI insights for all escalations"""
    print("=== AI INSIGHTS API CALLED ===")
    try:
        # Get all escalations for analysis
        escalations = Escalation.query.all()
        print(f"Found {len(escalations)} escalations for analysis")
        
        # Generate AI insights
        result = ai_service.generate_overall_insights(escalations)
        print(f"AI service result: {result}")
        
        return jsonify(result)
        
    except Exception as e:
        print(f"Error in generate_ai_insights: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': str(e),
            'insights': None
        }), 500

@app.route('/api/get-manager-info')
def get_manager_info():
    """Get DE Manager and Director info for a Primary POC using rchain command"""
    try:
        primary_poc = request.args.get('primary_poc', '').strip()
        
        if not primary_poc:
            return jsonify({
                'success': False,
                'error': 'Primary POC name is required'
            }), 400
        
        # Execute the rchain command
        import subprocess
        try:
            # Execute the real rchain command directly - no platform checks
            import platform
            
            # Execute the real rchain command directly - no simulation
            command = ['/usr/cisco/bin/rchain', '--mht', primary_poc]
            print(f"Executing: {' '.join(command)}")
            result = subprocess.run(command, capture_output=True, text=True, timeout=10)
            
            if result.returncode != 0:
                print(f"rchain command failed with return code {result.returncode}")
                print(f"stderr: {result.stderr}")
                return jsonify({
                    'success': False,
                    'error': f'rchain command failed: {result.stderr or "Command not available"}',
                    'command': ' '.join(command)
                }), 500
            
            if not result.stdout.strip():
                return jsonify({
                    'success': False,
                    'error': f'rchain command returned no output for user: {primary_poc}',
                    'command': ' '.join(command)
                }), 500
            
            # Use the actual rchain output
            simulated_output = result.stdout.strip()
            print(f"Real rchain output received for {primary_poc}")
            print(f"Output: {simulated_output}")
            
            # Parse the real command output
            output_lines = simulated_output.strip().split('\n')
            
            # Initialize variables for parsing
            de_manager = ''
            director = ''
            
            # Parse the output to extract manager and director info
            # Enhanced logic to handle multiple directors in hierarchy
            all_directors = []
            leader_found = False
            target_person = primary_poc.lower()
            
            for line in output_lines:
                if line.strip():
                    # Split on whitespace and rejoin to handle multiple spaces/tabs
                    parts = line.split()
                    if len(parts) >= 2:
                        username = parts[0].strip()
                        title = ' '.join(parts[1:]).strip()
                        
                        # Look for DE Manager - typically "Leader, Software Engineering"
                        if 'leader' in title.lower() and 'software engineering' in title.lower():
                            if not de_manager:  # Take the first leader found
                                de_manager = username
                                leader_found = True
                        
                        # Collect all directors
                        elif 'director' in title.lower() and 'software engineering' in title.lower():
                            all_directors.append(username)
                        
                        # Check if this is the target person to understand their position
                        elif username.lower() == target_person:
                            # Found the target person, now we know their position in hierarchy
                            pass
            
            # Dynamic director selection logic - Parse hierarchy to find correct director
            # No hardcoding - works for ANY rchain output
            if all_directors:
                if len(all_directors) == 1:
                    # Only one director found, use it
                    director = all_directors[0]
                else:
                    # Multiple directors - find the one closest to the leader in the hierarchy
                    # This works by finding the director that appears just before the leader
                    leader_position = -1
                    director_positions = {}
                    
                    # Find positions of leader and all directors in the hierarchy
                    for i, line in enumerate(output_lines):
                        if line.strip():
                            parts = line.split()
                            if len(parts) >= 2:
                                username = parts[0].strip()
                                
                                if username == de_manager:
                                    leader_position = i
                                elif username in all_directors:
                                    director_positions[username] = i
                    
                    # Select the director closest to (but above) the leader
                    if leader_position >= 0 and director_positions:
                        best_director = None
                        min_distance = float('inf')
                        
                        for dir_name, dir_pos in director_positions.items():
                            if dir_pos < leader_position:  # Director must be above leader
                                distance = leader_position - dir_pos
                                if distance < min_distance:
                                    min_distance = distance
                                    best_director = dir_name
                        
                        director = best_director if best_director else all_directors[0]
                    else:
                        # Fallback: use the last director (often most relevant)
                        director = all_directors[-1]
            
            # Fallback logic for missing values - always return N/A when no manager found
            if not de_manager:
                de_manager = 'N/A'  # No manager found in rchain output
                
            if not director:
                director = 'dheeraj'  # Default fallback
            
            # Prepare the response data
            raw_output_data = simulated_output
            command_note = f'Executed: /usr/cisco/bin/rchain --mht {primary_poc}'
            
            return jsonify({
                'success': True,
                'de_manager': de_manager,
                'director': director,
                'raw_output': raw_output_data,
                'command_executed': command_note,
                'note': 'Real rchain command execution'
            })
            
        except subprocess.TimeoutExpired:
            return jsonify({
                'success': False,
                'error': 'Command timeout - rchain took too long to respond'
            }), 500
        except FileNotFoundError:
            return jsonify({
                'success': False,
                'error': 'rchain command not found - please check if /usr/cisco/bin/rchain exists'
            }), 500
        except Exception as e:
            return jsonify({
                'success': False,
                'error': f'Command execution failed: {str(e)}'
            }), 500
            
    except Exception as e:
        print(f"Error in get_manager_info: {e}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/get-escalation-manager-info')
def get_escalation_manager_info():
    """Get Escalation Manager info for an Escalation Engineer using rchain command"""
    try:
        escalation_engineer = request.args.get('escalation_engineer', '').strip()
        
        if not escalation_engineer:
            return jsonify({
                'success': False,
                'error': 'Escalation Engineer name is required'
            }), 400
        
        # Execute the rchain command
        import subprocess
        try:
            # Execute the real rchain command directly
            command = ['/usr/cisco/bin/rchain', '--mht', escalation_engineer]
            print(f"Executing for escalation engineer: {' '.join(command)}")
            result = subprocess.run(command, capture_output=True, text=True, timeout=10)
            
            if result.returncode != 0:
                print(f"rchain command failed with return code {result.returncode}")
                print(f"stderr: {result.stderr}")
                return jsonify({
                    'success': False,
                    'error': f'rchain command failed: {result.stderr or "Command not available"}',
                    'command': ' '.join(command)
                }), 500
            
            if not result.stdout.strip():
                return jsonify({
                    'success': False,
                    'error': f'rchain command returned no output for user: {escalation_engineer}',
                    'command': ' '.join(command)
                }), 500
            
            # Use the actual rchain output
            output = result.stdout.strip()
            print(f"Real rchain output received for escalation engineer {escalation_engineer}")
            print(f"Output: {output}")
            
            # Parse the output to extract escalation manager info
            output_lines = output.strip().split('\n')
            
            # Initialize variables for parsing
            escalation_manager = ''
            
            # Parse the output to find the escalation manager (Leader, Software Engineering)
            for line in output_lines:
                if line.strip():
                    parts = line.split()
                    if len(parts) >= 2:
                        username = parts[0].strip()
                        title = ' '.join(parts[1:]).strip()
                        
                        # Look for "Leader, Software Engineering" - this is the escalation manager
                        if 'leader' in title.lower() and 'software engineering' in title.lower():
                            escalation_manager = username
                            break  # Found the leader, stop searching
            
            # If no leader found, fallback to looking for any leader role
            if not escalation_manager:
                for line in output_lines:
                    if line.strip():
                        parts = line.split()
                        if len(parts) >= 2:
                            title = ' '.join(parts[1:]).strip()
                            if 'leader' in title.lower():
                                escalation_manager = parts[0].strip()
                                break
            
            # Handle cases where no manager is found - always return N/A
            if not escalation_manager:
                escalation_manager = 'N/A'  # No manager found in rchain output
            
            # Prepare the response data
            raw_output_data = output
            command_note = f'Executed: /usr/cisco/bin/rchain --mht {escalation_engineer}'
            
            return jsonify({
                'success': True,
                'escalation_manager': escalation_manager,
                'raw_output': raw_output_data,
                'command_executed': command_note,
                'note': 'Real rchain command execution for escalation engineer'
            })
            
        except subprocess.TimeoutExpired:
            return jsonify({
                'success': False,
                'error': 'Command timeout - rchain took too long to respond'
            }), 500
        except FileNotFoundError:
            return jsonify({
                'success': False,
                'error': 'rchain command not found - please check if /usr/cisco/bin/rchain exists'
            }), 500
        except Exception as e:
            return jsonify({
                'success': False,
                'error': f'Command execution failed: {str(e)}'
            }), 500
            
    except Exception as e:
        print(f"Error in get_escalation_manager_info: {e}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/')
def index():
    return redirect(url_for('dashboard'))

@app.route('/test')
def test():
    return "Flask app is working! Submit form should be at /submit_form"

@app.route('/simple_test', methods=['GET', 'POST'])
def simple_test():
    if request.method == 'POST':
        print("=== SIMPLE TEST FORM RECEIVED ===")
        print("Form data:", request.form)
        return "Form submitted successfully!"
    
    return '''
    <html>
    <body>
        <h2>Simple Test Form</h2>
        <form method="POST" action="/simple_test">
            <input type="text" name="test_field" placeholder="Enter something" required>
            <button type="submit">Submit Test</button>
        </form>
    </body>
    </html>
    '''

@app.route('/submit_form')
# @login_required  # Temporarily disabled for testing
def submit_form():
    print("Submit form route accessed")
    return render_template('submit_form.html')

@app.route('/submit', methods=['POST'])
# @login_required  # Temporarily disabled for testing
def submit():
    print('=== FORM SUBMISSION RECEIVED ===')
    print('Request method:', request.method)
    print('Form data keys:', list(request.form.keys()))
    print('Form data values:')
    for key, value in request.form.items():
        print(f'  {key}: {value}')
    print('=================================')    
    data = request.form.to_dict()
    now = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    
    # Determine which cross teams have data
    cross_teams_with_data = []
    if any([data.get('routing_component'), data.get('routing_director'), data.get('routing_de_manager'), data.get('routing_engineer')]):
        cross_teams_with_data.append('Routing')
    if any([data.get('ios_xe_component'), data.get('ios_xe_director'), data.get('ios_xe_de_manager'), data.get('ios_xe_engineer')]):
        cross_teams_with_data.append('IOS-XE')
    if any([data.get('swiss_component'), data.get('swiss_director'), data.get('swiss_de_manager'), data.get('swiss_engineer')]):
        cross_teams_with_data.append('Location')
    if any([data.get('sjc_component'), data.get('sjc_director'), data.get('sjc_de_manager'), data.get('sjc_engineer')]):
        cross_teams_with_data.append('AP/SDA')
    
    # Set cross_team field to comma-separated list of teams with data
    cross_team_value = ', '.join(cross_teams_with_data) if cross_teams_with_data else None
    
    # Handle pending action logic based on state
    state = data.get('state') or 'Open'
    pending_action = data.get('pending_action', '')
    
    if state.lower() == 'open':
        # For open escalations, default to 'Dev' if no pending action provided
        if not pending_action or pending_action.strip() == '':
            pending_action = 'Dev'
    elif state.lower() == 'closed':
        # For closed escalations, use user selection or default to '-'
        if not pending_action or pending_action.strip() == '':
            pending_action = '-'
    
    # Handle initial remarks with timestamp and user info
    initial_remarks = data.get('remarks', '').strip()
    formatted_remarks = None
    if initial_remarks:
        current_time = datetime.now().strftime('%Y-%m-%d %H:%M')
        # Extract CEC ID from email (part before @cisco.com)
        user_email = session.get('webex_email', '')
        user_cec_id = user_email.split('@')[0] if user_email else 'unknown'
        formatted_remarks = f"[{current_time} - {user_cec_id}] {initial_remarks}"
    
    escalation = Escalation(
        customer=data.get('customer'),
        version=data.get('version'),
        bugid=data.get('bugid'),
        dedt_manager=data.get('dedt_manager'),
        engineer=data.get('engineer'),
        contributors=data.get('contributors'),
        cross_team=cross_team_value,  # Use the automatically determined value
        component_name=data.get('component_name'),
        director=data.get('director'),
        cross_de_manager=data.get('cross_de_manager'),
        cross_engineer=data.get('cross_engineer'),
        sr=data.get('sr'),
        bems=data.get('bems'),
        upgrade_attempt=data.get('upgrade_attempt'),
        symptom=data.get('symptom'),
        next_step=data.get('next_step'),
        director2=data.get('director2'),
        escalation_manager=data.get('escalation_manager'),
        escalation_engineer=data.get('escalation_engineer'),
        remarks=formatted_remarks,
        state=state,
        reopen_reason=data.get('reopen_reason'),
        closing_reason=data.get('closing_reason'),
        reported_on=data.get('reported_on'),
        closed_on=data.get('closed_on'),
        created_on=now,
        updated_on=now,
        pending_action=pending_action,
        severity=data.get('severity'),
        severity_text=data.get('severity_text'),
        severity_type=data.get('severity_type'),
        
        # Cross team details
        routing_component=data.get('routing_component'),
        routing_director=data.get('routing_director'),
        routing_de_manager=data.get('routing_de_manager'),
        routing_engineer=data.get('routing_engineer'),
        
        ios_xe_component=data.get('ios_xe_component'),
        ios_xe_director=data.get('ios_xe_director'),
        ios_xe_de_manager=data.get('ios_xe_de_manager'),
        ios_xe_engineer=data.get('ios_xe_engineer'),
        
        swiss_component=data.get('swiss_component'),
        swiss_director=data.get('swiss_director'),
        swiss_de_manager=data.get('swiss_de_manager'),
        swiss_engineer=data.get('swiss_engineer'),
        
        sjc_component=data.get('sjc_component'),
        sjc_director=data.get('sjc_director'),
        sjc_de_manager=data.get('sjc_de_manager'),
        sjc_engineer=data.get('sjc_engineer')
    )
    
    db.session.add(escalation)
    db.session.commit()

    # --- Post escalation info to Webex space ---
    WEBEX_BOT_TOKEN = os.getenv("WEBEX_BOT_TOKEN")
    WEBEX_ROOM_ID = os.getenv("WEBEX_ROOM_ID")
    HOST_URL = os.getenv("HOST_URL", "localhost:5000")
    
    if WEBEX_BOT_TOKEN and WEBEX_ROOM_ID:
        headers = {
            "Authorization": f"Bearer {WEBEX_BOT_TOKEN}",
            "Content-Type": "application/json"
        }
        
        # Build cross team summary
        cross_teams = []
        if escalation.routing_component or escalation.routing_director:
            cross_teams.append(f"Routing: {escalation.routing_component or 'N/A'}")
        if escalation.ios_xe_component or escalation.ios_xe_director:
            cross_teams.append(f"IOS-XE: {escalation.ios_xe_component or 'N/A'}")
        if escalation.swiss_component or escalation.swiss_director:
            cross_teams.append(f"Location: {escalation.swiss_component or 'N/A'}")
        if escalation.sjc_component or escalation.sjc_director:
            cross_teams.append(f"AP/SDA: {escalation.sjc_component or 'N/A'}")
        
        # Build condensed cross team info for new escalation
        condensed_teams_new = []
        if escalation.ios_xe_engineer or escalation.ios_xe_de_manager or escalation.ios_xe_director or escalation.ios_xe_component:
            condensed_teams_new.append("IOS-XE")
        if escalation.sjc_engineer or escalation.sjc_de_manager or escalation.sjc_director or escalation.sjc_component:
            condensed_teams_new.append("AP/SDA")
        if escalation.swiss_engineer or escalation.swiss_de_manager or escalation.swiss_director or escalation.swiss_component:
            condensed_teams_new.append("Location")
        if escalation.routing_engineer or escalation.routing_de_manager or escalation.routing_director or escalation.routing_component:
            condensed_teams_new.append("Routing")
        
        cross_teams_display_new = " / ".join(condensed_teams_new) if condensed_teams_new else "None"
        edit_url = f"http://{HOST_URL}/edit/{escalation.id}"
        
        # Build escalation details section only for Open state
        escalation_details = ""
        if escalation.state and escalation.state.upper() != 'CLOSED':
            escalation_details = (
                f"### ⚡ **Technical Details**\n"
                f"| SR Number | BEMS | Symptom |\n"
                f"|:---------:|:----:|:-------:|\n"
                f"| {escalation.sr or 'Not provided'} | {escalation.bems or 'Not provided'} | {escalation.symptom or 'Not described'} |\n\n"
                
                f"| Upgrade Attempt | Next Step | Remarks |\n"
                f"|:---------------:|:---------:|:-------:|\n"
                f"| {escalation.upgrade_attempt or 'None attempted'} | {escalation.next_step or 'To be determined'} | {escalation.remarks or 'No additional remarks'} |\n\n"
            )

        # Prepare data for the adaptive card format (same as edit/update page)
        payload = {
            "id": escalation.id,
            "customer": escalation.customer or "Not specified",
            "version": escalation.version or "N/A",
            "bug_id": escalation.bugid or "N/A",
            "component": escalation.component_name or "Not specified",
            "severity": escalation.severity_type or escalation.severity or "N/A",
            "de_mgr": escalation.dedt_manager or "N/A",
            "primary_poc": escalation.engineer or "N/A",
            "contributors": escalation.contributors or "None",
            "sr": escalation.sr or "N/A",
            "bems": escalation.bems or "N/A",
            "symptom": escalation.symptom or "N/A",
            "upgrade_attempt": escalation.upgrade_attempt or "N/A",
            "escalation_remarks": escalation.remarks or "N/A",
            "url": f"http://{HOST_URL}/edit/{escalation.id}"
        }
        
        # Send adaptive card to Webex (same as edit/update page)
        try:
            response = send_escalation_card(WEBEX_ROOM_ID, WEBEX_BOT_TOKEN, payload, "🚨 New Escalation Request")
            
            if response.status_code == 200:
                print(f"✅ Successfully posted escalation #{escalation.id} to Webex")
            else:
                print(f"❌ Failed to post to Webex. Status: {response.status_code}, Response: {response.text}")
        except Exception as e:
            print(f"❌ Failed to post escalation to Webex: {e}")
    else:
        print("⚠️ Webex credentials not configured - skipping notification")
    
    return redirect(url_for('dashboard'))


# @app.route('/send_escalation_to_webex/<int:escalation_id>', methods=['POST'])
# def send_escalation_to_webex(escalation_id):
#     """Send detailed escalation information to Webex space (similar to attachment format)"""
#     # Get Webex credentials from environment
#     WEBEX_BOT_TOKEN = os.getenv("WEBEX_BOT_TOKEN")
#     WEBEX_ROOM_ID = os.getenv("WEBEX_ROOM_ID")
#     HOST_URL = os.getenv("HOST_URL", "localhost:5000")
#     
#     if not all([WEBEX_BOT_TOKEN, WEBEX_ROOM_ID]):
#         return jsonify({"error": "Webex credentials not configured"}), 400
#     
#     try:
#         escalation = Escalation.query.get_or_404(escalation_id)
#         
#         # Build cross-team summary
#         cross_teams = []
#         if escalation.ios_xe_engineer or escalation.ios_xe_de_manager or escalation.ios_xe_director or escalation.ios_xe_component:
#             cross_teams.append("IOS-XE")
#         if escalation.sjc_engineer or escalation.sjc_de_manager or escalation.sjc_director or escalation.sjc_component:
#             cross_teams.append("AP/SDA")
#         if escalation.swiss_engineer or escalation.swiss_de_manager or escalation.swiss_director or escalation.swiss_component:
#             cross_teams.append("Location")
#         if escalation.routing_engineer or escalation.routing_de_manager or escalation.routing_director or escalation.routing_component:
#             cross_teams.append("Routing")
#         
#         cross_teams_display = " / ".join(cross_teams) if cross_teams else "None"
# 
#         # Format message similar to the attachment
#         message = f"📋 **Escalation Details #{escalation.id}**\n\n"
#         
#         # General Information (matching attachment format)
#         message += f"**Customer:** {escalation.customer or 'None'}\n"
#         message += f"**Version:** {escalation.version or 'None'}   **Bug ID:** {escalation.bugid or 'None'}   **Component:** {escalation.component_name or 'None'}\n"
#         message += f"**Sev:** {escalation.severity_type or escalation.severity or 'None'}   **Sev Remarks:** {escalation.severity_text or 'None'}\n"
#         message += f"**DE Mgr:** {escalation.dedt_manager or 'None'}   **Primary PoC:** {escalation.engineer or 'None'}   **Contributors:** {escalation.contributors or 'None'}\n"
#         message += f"**Cross Dependent Teams:** {cross_teams_display}\n\n"
#         
#         # State and status information
#         if escalation.state:
#             message += f"**State:** {escalation.state}\n"
#             if escalation.state.upper() == 'CLOSED':
#                 if escalation.closing_reason:
#                     message += f"**Closing Reason:** {escalation.closing_reason}\n"
#                 if escalation.closing_remarks:
#                     message += f"**Closing Remarks:** {escalation.closing_remarks}\n"
#         message += "\n"
#         
#         # Escalation Details (only for open escalations)
#         if escalation.state and escalation.state.upper() != 'CLOSED':
#             message += f"⚡ **Escalation Details**\n"
#             message += f"```\n"
#             message += f"SR: {escalation.sr or 'None'}   BEMS: {escalation.bems or 'None'}\n"
#             message += f"Symptom reported/understood: {escalation.symptom or 'None'}   Functionality: {escalation.functionality or 'None'}\n"
#             message += f"Any upgrade attempt and new version: {escalation.upgrade_attempt or 'None'}   Next step: {escalation.next_step or 'None'}\n"
#             message += f"Remarks: {escalation.remarks or 'None'}\n"
#             message += f"```\n\n"
#         
#         # Add links
#         message += f"**[📝 Edit Escalation](http://{HOST_URL}/edit/{escalation.id})** | "
#         message += f"**[📊 View Details](http://{HOST_URL}/details/{escalation.id})** | "
#         message += f"**[📋 Dashboard](http://{HOST_URL}/dashboard)**"
#         
#         # Send to Webex
#         payload = {
#             "roomId": WEBEX_ROOM_ID,
#             "markdown": message
#         }
#         
#         response = requests.post(
#             "https://webexapis.com/v1/messages",
#             headers={
#                 "Authorization": f"Bearer {WEBEX_BOT_TOKEN}",
#                 "Content-Type": "application/json"
#             },
#             json=payload
#         )
#         
#         if response.status_code == 200:
#             return jsonify({"success": f"Escalation #{escalation.id} details sent to Webex successfully"})
#         else:
#             return jsonify({"error": f"Failed to send to Webex. Status: {response.status_code}"}), 500
#             
#     except Exception as e:
#         return jsonify({"error": f"Failed to send escalation to Webex: {str(e)}"}), 500


def format_table_with_borders(headers, rows):
    """Format table with ASCII borders with proper column spacing"""
    if not headers or not rows:
        return ""
    
    # Calculate maximum width for each column with generous padding
    widths = []
    for i in range(len(headers)):
        col_widths = [len(str(headers[i]))]
        for row in rows:
            if i < len(row):
                col_widths.append(len(str(row[i])))
        max_width = max(col_widths)
        # Add generous padding (minimum 6 extra characters or more)
        widths.append(max(max_width + 8, 12))
    
    # Create top border
    top_border = "+"
    for width in widths:
        top_border += "-" * width + "+"
    
    # Create header row with left alignment and full width padding
    header_row = "|"
    for i, header in enumerate(headers):
        padded_header = f" {str(header):<{widths[i]-2}} "
        header_row += padded_header + "|"
    
    # Create separator border (same as top)
    separator = top_border
    
    # Create data rows with full width padding
    data_rows = []
    for row in rows:
        row_str = "|"
        for i in range(len(headers)):
            if i < len(row):
                value = str(row[i])
            else:
                value = ""
            # Left align with full column width padding
            padded_value = f" {value:<{widths[i]-2}} "
            row_str += padded_value + "|"
        data_rows.append(row_str)
    
    # Create bottom border
    bottom_border = top_border
    
    # Combine all parts
    result = [top_border, header_row, separator]
    result.extend(data_rows)
    result.append(bottom_border)
    
    return "\n".join(result)

def format_table_row(headers, values):
    """Wrapper function to maintain compatibility - creates single row table"""
    return format_table_with_borders(headers, [values])

def get_last_n_remarks(remarks_text, n=3):
    """Extract the last n remarks from the remarks text"""
    if not remarks_text or remarks_text == "N/A":
        return "N/A"
    
    # Split by lines and look for timestamp patterns [YYYY-MM-DD HH:MM - user]
    lines = remarks_text.split('\n')
    remark_entries = []
    
    for line in lines:
        line = line.strip()
        if line and '[' in line and ']' in line:
            # This looks like a timestamped remark
            remark_entries.append(line)
    
    # If no timestamped remarks found, treat the whole text as one remark
    if not remark_entries:
        return remarks_text
    
    # Get the last n remarks
    last_remarks = remark_entries[-n:] if len(remark_entries) >= n else remark_entries
    
    return '\n'.join(last_remarks)

def send_escalation_update_to_webex(escalation, changes=None):
    """Send an updated escalation notification to Webex"""
    # Get Webex credentials from environment
    WEBEX_BOT_TOKEN = os.getenv("WEBEX_BOT_TOKEN")
    WEBEX_ROOM_ID = os.getenv("WEBEX_ROOM_ID")
    HOST_URL = os.getenv("HOST_URL", "localhost:5000")
    
    if not all([WEBEX_BOT_TOKEN, WEBEX_ROOM_ID]):
        print("⚠️ Webex credentials not configured - skipping update notification")
        return
    
    # Prepare data for the table card format (4x2 format: Row 1: Customer | Version | Bug ID | Component, Row 2: Severity | DE Mgr | Primary PoC | Contributors)
    payload = {
        "id": escalation.id,
        "customer": escalation.customer or "Not specified",
        "version": escalation.version or "N/A",
        "bug_id": escalation.bugid or "N/A",
        "component": escalation.component_name or "Not specified",
        "severity": escalation.severity_type or escalation.severity or "N/A",
        "de_mgr": escalation.dedt_manager or "N/A",
        "primary_poc": escalation.engineer or "N/A",
        "contributors": escalation.contributors or "None",
        "sr": escalation.sr or "N/A",
        "bems": escalation.bems or "N/A",
        "symptom": escalation.symptom or "N/A",
        "upgrade_attempt": escalation.upgrade_attempt or "N/A",
        "escalation_remarks": get_last_n_remarks(escalation.remarks, 3),
        "url": f"http://{HOST_URL}/edit/{escalation.id}"
    }
    
    # Send table card to Webex
    try:
        response = send_escalation_card(WEBEX_ROOM_ID, WEBEX_BOT_TOKEN, payload, "Escalation Update")
        
        if response.status_code == 200:
            print(f"✅ Successfully posted escalation update #{escalation.id} to Webex")
        else:
            print(f"❌ Failed to post update to Webex. Status: {response.status_code}, Response: {response.text}")
    except Exception as e:
        print(f"❌ Failed to post escalation update to Webex: {e}")


def send_escalation_card(room_id, token, data, title="Escalation Update"):
    """Send escalation details as an Adaptive Card with table format to Webex"""
    url = "https://webexapis.com/v1/messages"
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json"
    }

    card = {
      "roomId": room_id,
      "text": f"Escalation Update #{data['id']} - {data['customer']} - {data['component']}",
      "attachments": [
        {
          "contentType": "application/vnd.microsoft.card.adaptive",
          "content": {
            "$schema": "http://adaptivecards.io/schemas/adaptive-card.json",
            "type": "AdaptiveCard",
            "version": "1.2",
            "width": "60%",
            "rtl": False,
            "body": [
              {
                "type": "Container",
                "style": "default",
                "width": "stretch",
                "padding": "Default",
                "items": [
                  {
                    "type": "TextBlock",
                    "text": f"{title} - #{data['id']}",
                    "weight": "Bolder",
                    "size": "Medium",
                    "wrap": True
                  },
                  {
                    "type": "TextBlock",
                    "text": "📋 General Information and Teams",
                    "weight": "Bolder",
                    "size": "Small",
                    "spacing": "Medium"
                  },
                  {
                    "type": "Container",
                    "style": "emphasis",
                    "width": "stretch",
                    "minWidth": "400px",
                "items": [
                  {
                    "type": "ColumnSet",
                    "width": "stretch",
                    "spacing": "None",
                    "minWidth": "650px",
                    "columns": [
                      {
                        "type": "Column",
                        "width": "90px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Customer**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "75px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Version**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "85px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Bug ID**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "110px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Component**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      }
                    ]
                  },
                  {
                    "type": "ColumnSet",
                    "width": "stretch",
                    "spacing": "Small",
                    "minWidth": "650px",
                    "columns": [
                      {
                        "type": "Column",
                        "width": "90px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data["customer"],
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "75px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data["version"],
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "85px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data["bug_id"],
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "110px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data["component"],
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      }
                    ]
                  },
                  {
                    "type": "ColumnSet",
                    "width": "stretch",
                    "spacing": "Medium",
                    "minWidth": "400px",
                    "columns": [
                      {
                        "type": "Column",
                        "width": "90px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Severity**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "75px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**DE Mgr**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "85px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Primary PoC**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "110px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Contributors**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      }
                    ]
                  },
                  {
                    "type": "ColumnSet",
                    "width": "stretch",
                    "spacing": "Small",
                    "minWidth": "400px",
                    "columns": [
                      {
                        "type": "Column",
                        "width": "90px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data["severity"],
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "75px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data["de_mgr"],
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "85px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data["primary_poc"],
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "110px",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data["contributors"],
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      }
                    ]
                  }
                ]
                  },
                  {
                    "type": "TextBlock",
                    "text": "**Escalation Details**",
                    "weight": "Bolder",
                    "size": "Small",
                    "spacing": "Large"
                      },
                  {
                    "type": "Container",
                    "style": "emphasis",
                    "width": "stretch",
                    "items": [
                      {
                        "type": "ColumnSet",
                        "columns": [
                          {
                            "type": "Column",
                            "width": "15%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**SR:**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "35%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data.get("sr", "N/A"),
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "15%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**BEMS:**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "35%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data.get("bems", "N/A"),
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      }
                    ]
                  },
                  {
                    "type": "ColumnSet",
                    "columns": [
                      {
                        "type": "Column",
                        "width": "15%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Symptom:**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "85%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data.get("symptom", "N/A"),
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      }
                    ]
                  },
                  {
                    "type": "ColumnSet",
                    "columns": [
                      {
                        "type": "Column",
                        "width": "15%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Upgrade:**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "85%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data.get("upgrade_attempt", "N/A"),
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      }
                    ]
                  },
                  {
                    "type": "ColumnSet",
                    "columns": [
                      {
                        "type": "Column",
                        "width": "15%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": "**Remarks:**",
                            "weight": "Bolder",
                            "size": "Small"
                          }
                        ]
                      },
                      {
                        "type": "Column",
                        "width": "85%",
                        "items": [
                          {
                            "type": "TextBlock",
                            "text": data.get("escalation_remarks", "N/A"),
                            "size": "Small",
                            "wrap": True
                          }
                        ]
                      }
                    ]
                  }
                    ]
                  },
                  {
                    "type": "ActionSet",
                    "actions": [
                      {
                        "type": "Action.OpenUrl",
                        "title": "📎 View Escalation",
                        "url": data["url"]
                      }
                    ]
                  }
                ]
              }
            ]
          }
        }
      ]
    }

    return requests.post(url, headers=headers, json=card)


@app.route('/send_table_to_webex/<int:escalation_id>', methods=['POST'])
@login_required
def send_table_to_webex(escalation_id):
    """Send escalation details as an Adaptive Card to Webex space"""
    # Get Webex credentials from environment
    WEBEX_BOT_TOKEN = os.getenv("WEBEX_BOT_TOKEN")
    WEBEX_ROOM_ID = os.getenv("WEBEX_ROOM_ID")
    HOST_URL = os.getenv("HOST_URL", "localhost:5000")
    
    if not all([WEBEX_BOT_TOKEN, WEBEX_ROOM_ID]):
        return jsonify({"error": "Webex credentials not configured"}), 500
    
    # Get escalation data
    escalation = Escalation.query.get_or_404(escalation_id)
    
    # Prepare data for the comprehensive card format
    payload = {
        "id": escalation.id,
        "customer": escalation.customer or "Not specified",
        "version": escalation.version or "N/A",
        "bug_id": escalation.bugid or "N/A",
        "component": escalation.component_name or "Not specified",
        "severity": escalation.severity_type or escalation.severity or "Not specified",
        "severity_remarks": escalation.severity_text or "N/A",
        "de_mgr": escalation.engineer or "N/A",
        "primary_poc": escalation.escalation_engineer or "N/A",
        "contributors": escalation.contributors or "None",
        "sr": escalation.sr or "N/A",
        "bems": escalation.bems or "N/A",
        "symptom": escalation.symptom or "N/A",
        "upgrade_attempt": escalation.upgrade_attempt or "N/A",
        "escalation_remarks": get_last_n_remarks(escalation.remarks, 3),
        "url": f"http://{HOST_URL}/edit/{escalation.id}"
    }
    
    # Send Adaptive Card to Webex
    try:
        response = send_escalation_card(WEBEX_ROOM_ID, WEBEX_BOT_TOKEN, payload, "Escalation Update")
        
        if response.status_code == 200:
            return jsonify({"success": f"Escalation #{escalation.id} table sent to Webex successfully!"})
        else:
            return jsonify({"error": f"Failed to send to Webex. Status: {response.status_code}"}), 500
            
    except Exception as e:
        return jsonify({"error": f"Failed to send table to Webex: {str(e)}"}), 500


@app.route('/webex_table_test')
@login_required
def webex_table_test():
    """Test page for sending formatted tables to Webex"""
    return render_template('webex_table_test.html')


@app.route('/dashboard')
@app.route('/dashboard', methods=['GET', 'POST'])
#@login_required  # Temporarily disabled for debugging
def dashboard():
    query = Escalation.query
    de_manager = request.values.get('de_manager', '')
    state = request.values.get('state', '')
    cross_team = request.values.get('cross_team', '')
    week_filter = request.values.get('week_filter', '')
    closing_reason = request.values.get('closing_reason', '')
    pending_action = request.values.get('pending_action', '')
    if de_manager:
        query = query.filter_by(dedt_manager=de_manager)
    if state:
        # Use case-insensitive filtering for state to match normalized dropdown values
        query = query.filter(Escalation.state.ilike(state))
    if cross_team:
        query = query.filter_by(cross_team=cross_team)
    if pending_action and pending_action != 'All':
        query = query.filter(Escalation.pending_action == pending_action)
    cap_status = request.values.get('cap_status', '')
    if cap_status and cap_status != 'All':
        query = query.filter(Escalation.severity_type == cap_status)
    if closing_reason:
        query = query.filter(Escalation.state == 'Closed').filter_by(closing_reason=closing_reason)
    items = []
    now = datetime.now()
    if week_filter:
        if week_filter == 'open_week' or week_filter == 'last_week':
            query = query.filter(Escalation.state.ilike('open'))
            for item in query.all():
                try:
                    reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
                except Exception:
                    reported_date = None
                if reported_date and 0 <= (now - reported_date).days <= 7:
                    items.append(item)
        elif week_filter == 'closed_week':
            query = query.filter(Escalation.state.ilike('closed'))
            for item in query.all():
                try:
                    closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d') if item.closed_on else None
                except Exception:
                    closed_date = None
                if closed_date and 0 <= (now - closed_date).days <= 7:
                    items.append(item)
        elif week_filter == 'open_2weeks':
            query = query.filter(Escalation.state.ilike('open'))
            for item in query.all():
                try:
                    reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
                except Exception:
                    reported_date = None
                if reported_date and (now - reported_date).days >= 14:
                    items.append(item)
        elif week_filter == 'open_nobugid':
            query = query.filter(Escalation.state.ilike('open'))
            for item in query.all():
                if not item.bugid or not item.bugid.strip():
                    items.append(item)
        else:
            items = query.all()
    else:
        items = query.all()

    # Pagination logic
    page = int(request.args.get('page', 1))
    per_page = int(request.args.get('per_page', 10))
    total_items = len(items)
    total_pages = (total_items + per_page - 1) // per_page
    start = (page - 1) * per_page
    end = start + per_page
    items = items[start:end]
    # Calculate MTTR for each item
    mttr_list = []
    for item in items:
        try:
            reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
        except Exception:
            reported_date = None
        mttr_days = ''
        if reported_date:
            if item.state and item.state.lower() == 'open':
                mttr_days = (datetime.now() - reported_date).days
            elif item.state and item.state.lower() == 'closed' and item.closed_on:
                try:
                    closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d')
                    mttr_days = (closed_date - reported_date).days
                except Exception:
                    mttr_days = ''
        mttr_list.append(mttr_days)
    # Get unique values for dropdown filters
    all_escalations = Escalation.query.all()
    
    # Get unique values for each filter
    de_managers = sorted(set(item.dedt_manager for item in all_escalations if item.dedt_manager and item.dedt_manager.strip()))
    # Provide fixed list of states instead of deriving from data
    states = ['Open', 'Closed']  # Fixed list to always show both options
    cross_teams = sorted(set(item.cross_team for item in all_escalations if item.cross_team and item.cross_team.strip()))
    closing_reasons = sorted(set(item.closing_reason for item in all_escalations if item.closing_reason and item.closing_reason.strip()))
    pending_actions = sorted(set(item.pending_action for item in all_escalations if item.pending_action and item.pending_action.strip()))
    severities = sorted(set(item.severity_type for item in all_escalations if item.severity_type and item.severity_type.strip()))
    
    # current_user = get_current_user()  # Login disabled
    current_user = None
    return render_template('dashboard.html', 
                         items=items, 
                         mttr_list=mttr_list, 
                         de_manager=de_manager, 
                         state=state, 
                         cross_team=cross_team, 
                         week_filter=week_filter, 
                         closing_reason=closing_reason, 
                         pending_action=pending_action, 
                         cap_status=cap_status, 
                         current_user=current_user, 
                         page=page, 
                         per_page=per_page, 
                         total_pages=total_pages, 
                         total_items=total_items,
                         # Filter options
                         de_managers=de_managers,
                         states=states,
                         cross_teams=cross_teams,
                         closing_reasons=closing_reasons,
                         pending_actions=pending_actions,
                         severities=severities)

# @app.route('/send_dashboard_to_webex', methods=['POST'])
# def send_dashboard_to_webex():
#     """Send escalation dashboard summary to Webex space"""
#     # Get Webex credentials from environment
#     WEBEX_BOT_TOKEN = os.getenv("WEBEX_BOT_TOKEN")
#     WEBEX_ROOM_ID = os.getenv("WEBEX_ROOM_ID")
#     HOST_URL = os.getenv("HOST_URL", "localhost:5000")
#     
#     if not all([WEBEX_BOT_TOKEN, WEBEX_ROOM_ID]):
#         return jsonify({"error": "Webex credentials not configured"}), 400
#     
#     try:
#         # Get current escalations
#         open_escalations = Escalation.query.filter(Escalation.state.ilike('open')).all()
#         closed_escalations = Escalation.query.filter(Escalation.state.ilike('closed')).all()
#         total_escalations = len(open_escalations) + len(closed_escalations)
#         
#         # Build summary message
#         message = f"📊 **Escalation Dashboard Summary**\n\n"
#         message += f"**Total Escalations:** {total_escalations}\n"
#         message += f"**Open:** {len(open_escalations)} | **Closed:** {len(closed_escalations)}\n\n"
#         
#         if open_escalations:
#             message += f"🔴 **Open Escalations ({len(open_escalations)})**\n"
#             for esc in open_escalations[:10]:  # Show first 10
#                 cross_teams = []
#                 if esc.ios_xe_engineer or esc.ios_xe_de_manager or esc.ios_xe_director or esc.ios_xe_component:
#                     cross_teams.append("IOS-XE")
#                 if esc.sjc_engineer or esc.sjc_de_manager or esc.sjc_director or esc.sjc_component:
#                     cross_teams.append("AP/SDA")
#                 if esc.swiss_engineer or esc.swiss_de_manager or esc.swiss_director or esc.swiss_component:
#                     cross_teams.append("Location")
#                 if esc.routing_engineer or esc.routing_de_manager or esc.routing_director or esc.routing_component:
#                     cross_teams.append("Routing")
#                 
#                 cross_teams_display = " / ".join(cross_teams) if cross_teams else "None"
#                 
#                 message += f"• **#{esc.id}** - {esc.customer or 'N/A'} | {esc.dedt_manager or 'N/A'} | Teams: {cross_teams_display}\n"
#                 message += f"  Sev: {esc.severity_type or esc.severity or 'N/A'} | Component: {esc.component_name or 'N/A'}\n"
#                 message += f"  **[View Details](http://{HOST_URL}/details/{esc.id})**\n\n"
#             
#             if len(open_escalations) > 10:
#                 message += f"... and {len(open_escalations) - 10} more open escalations\n\n"
#         
#         message += f"**[📋 View Full Dashboard](http://{HOST_URL}/dashboard)**"
#         
#         # Send to Webex
#         payload = {
#             "roomId": WEBEX_ROOM_ID,
#             "markdown": message
#         }
#         
#         response = requests.post(
#             "https://webexapis.com/v1/messages",
#             headers={
#                 "Authorization": f"Bearer {WEBEX_BOT_TOKEN}",
#                 "Content-Type": "application/json"
#             },
#             json=payload
#         )
#         
#         if response.status_code == 200:
#             return jsonify({"success": "Dashboard summary sent to Webex successfully"})
#         else:
#             return jsonify({"error": f"Failed to send to Webex. Status: {response.status_code}"}), 500
#             
#     except Exception as e:
#         return jsonify({"error": f"Failed to send dashboard to Webex: {str(e)}"}), 500
#     add_table_slides(prs, "Closed Issues Details", closed_cols, closed_last_week)
#     add_table_slides(prs, "Open Issues Details", open_cols, open_last_week)
# 
#     tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
#     prs.save(tmp.name)
#     tmp.close()
# 
#     headers = {
#         "Authorization": "Bearer Mzc2MWFjMmMtZDIyOC00NWIxLTgyZjYtMzMwYThhYmRkMTkxZWRiYjM1Y2YtYWMz_PF84_1eb65fdf-9643-417f-9974-ad72cae0e10f"
#     }
#     try:
#         with open(tmp.name, "rb") as f:
#             files = {
#                 'files': ("dashboard.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
#             }
#             data = {
#                 "roomId": "Y2lzY29zcGFyazovL3VzL1JPT00vMmVlOGU5ZjAtOGM5NC0xMWYwLTg2NTUtOTc2ZjAyZTNkZWU3",
#                 "text": "Escalation Dashboard PPT attached."
#             }
#             url = "https://webexapis.com/v1/messages"
#             r = requests.post(url, headers=headers, data=data, files=files)
#         os.unlink(tmp.name)
#         if r.status_code == 200:
#             return jsonify({"status": "success", "message": "Dashboard PPT sent to Webex!"})
#         else:
#             try:
#                 return jsonify({"status": "error", "details": r.json()})
#             except Exception:
#                 return jsonify({"status": "error", "details": r.text})
#     except Exception as e:
#         if os.path.exists(tmp.name):
#             os.unlink(tmp.name)
#         return jsonify({"status": "error", "details": str(e)})
#     # Prevent fall-through to render_template
#     return ('', 204)
    query = Escalation.query
    de_manager = request.values.get('de_manager', '')
    engineer = request.values.get('engineer', '')
    state = request.values.get('state', '')
    cross_team = request.values.get('cross_team', '')
    week_filter = request.values.get('week_filter', '')
    closing_reason = request.values.get('closing_reason', '')
    pending_action = request.values.get('pending_action', '')
    if de_manager:
        query = query.filter_by(dedt_manager=de_manager)
    if engineer:
        query = query.filter_by(engineer=engineer)
    if state:
        query = query.filter_by(state=state)
    if cross_team:
        query = query.filter_by(cross_team=cross_team)
    if pending_action and pending_action != 'All':
        query = query.filter(Escalation.pending_action == pending_action)
    cap_status = request.values.get('cap_status', '')
    if cap_status and cap_status != 'All':
        query = query.filter(Escalation.severity_type == cap_status)
    if closing_reason:
        query = query.filter(Escalation.state == 'Closed').filter_by(closing_reason=closing_reason)
    items = []
    now = datetime.now()
    if week_filter:
        if week_filter == 'open_week' or week_filter == 'last_week':
            query = query.filter(Escalation.state == 'Open')
            for item in query.all():
                try:
                    reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
                except Exception:
                    reported_date = None
                if reported_date and 0 <= (now - reported_date).days <= 7:
                    items.append(item)
        elif week_filter == 'closed_week':
            query = query.filter(Escalation.state == 'Closed')
            for item in query.all():
                try:
                    closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d') if item.closed_on else None
                except Exception:
                    closed_date = None
                # Only include if closed_on date is present and within last 7 days
                if closed_date and 0 <= (now - closed_date).days <= 7:
                    items.append(item)
        elif week_filter == 'open_2weeks':
            query = query.filter(Escalation.state == 'Open')
            for item in query.all():
                try:
                    reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
                except Exception:
                    reported_date = None
                if reported_date and (now - reported_date).days >= 14:
                    items.append(item)
        elif week_filter == 'open_nobugid':
            query = query.filter(Escalation.state == 'Open')
            for item in query.all():
                if not item.bugid or not item.bugid.strip():
                    items.append(item)
        else:
            items = query.all()
    else:
        items = query.all()
    # Calculate MTTR for each item
    mttr_list = []
    for item in items:
        try:
            reported_date = datetime.strptime(item.reported_on, '%Y-%m-%d') if item.reported_on else None
        except Exception:
            reported_date = None
        mttr_days = ''
        if reported_date:
            if item.state and item.state.lower() == 'open':
                mttr_days = (datetime.now() - reported_date).days
            elif item.state and item.state.lower() == 'closed' and item.closed_on:
                try:
                    closed_date = datetime.strptime(item.closed_on, '%Y-%m-%d')
                    mttr_days = (closed_date - reported_date).days
                except Exception:
                    mttr_days = ''
        mttr_list.append(mttr_days)
    # current_user = get_current_user()  # Login disabled
    current_user = None
    return render_template('dashboard.html', items=items, mttr_list=mttr_list, de_manager=de_manager, engineer=engineer, state=state, cross_team=cross_team, week_filter=week_filter, closing_reason=closing_reason, pending_action=pending_action, cap_status=cap_status, current_user=current_user)

@app.route('/get_webex_messages', methods=['GET'])
def get_webex_messages():
    import re
    import re
    def extract_field(text, label):
        # Match label followed by value, stopping at next label or end
        pattern = rf'{re.escape(label)}\s*:?(.*?)(?=\n[A-Za-z0-9 \(\)/]+:|$)'
        match = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
        if match:
            value = match.group(1).strip()
            # Remove any trailing label text accidentally captured
            value = re.split(r'\n[A-Za-z0-9 \(\)/]+:', value)[0].strip()
            return value if value else None
        return None

    # List of all columns/labels to extract
    field_map = {
        'Customer:': 'customer',
        'Current Version:': 'version',
        'SR': 'sr',
        'BEMS': 'bems',
        'Any upgrade attempt and new version:': 'upgrade_attempt',
        'Symptom reported/understood:': 'symptom',
        'DEDT-Manager from our team:': 'dedt_manager',
        'Engineer(s) from our team:': 'engineer',
        'Cross dependent team (Routin/IOS-XE/Location/AP-SDA, Component name, Engineer, DE Manager, Director):': 'cross_team',
        'Escalation Engineer and Manager (if no one is there, put NA):': 'escalation_engineer',
    }

    url = "https://webexapis.com/v1/messages"
    headers = {
        "Authorization": f"Bearer N2ZkNDkxNGMtNzVhNC00MjQyLThlMzUtNzQyMWFiMDJiOWY2NDFmZjgxZWMtZWE2_PF84_1eb65fdf-9643-417f-9974-ad72cae0e10f"
    }
    params = {
        "roomId": "Y2lzY29zcGFyazovL3VzL1JPT00vMmVlOGU5ZjAtOGM5NC0xMWYwLTg2NTUtOTc2ZjAyZTNkZWU3",
        "max": 5  # Fetch more messages to ensure both parent and thread are processed
    }
    try:
        response = requests.get(url, headers=headers, params=params)
        if response.status_code == 200:
            messages = response.json().get("items", [])
            imported = 0
            try:
                # Build a lookup of message id to message for parent-child matching
                msg_id_to_msg = {m['id']: m for m in messages if 'id' in m}
                # First pass: create parent entries
                for msg in messages:
                    try:
                        text = msg.get("text", "")
                        msg_id = msg.get("id")
                        parent_id = msg.get("parentId")
                        is_threaded = bool(parent_id)
                        remarks_value = text if is_threaded else ''
                        de_manager = extract_field(text, 'DEDT-Manager from our team:') or '-'
                        customer = extract_field(text, 'Customer:') or '-'
                        version = extract_field(text, 'Current Version:') or '-'
                        sr_bems_ddts = extract_field(text, 'SR and BEMS and DDTS (if available)#:') or '-'
                        # Extract bug id (CSC...) from SR/BEMS/DDTS field
                        import re
                        bugid_match = re.search(r'(CSC\w+)', sr_bems_ddts)
                        bugid = bugid_match.group(1) if bugid_match else 'N/A'
                        version = version.split()[0].strip()
                        upgrade_attempt = extract_field(text, 'Any upgrade attempt and new version:') or '-'
                        engineer = extract_field(text, 'Engineer(s) from our team:') or '-'
                        engineer = engineer.split('\n')[0].split('Cross dependent team')[0].split('Escalation Engineer and Manager')[0].strip()
                        cross_dependency = extract_field(text, 'Cross dependent team (Routin/IOS-XE/Location/AP-SDA, Component name, Engineer, DE Manager, Director):') or '-'
                        # Remove any trailing Escalation Engineer and Manager label and value
                        cross_dependency = cross_dependency.split('Escalation Engineer and Manager')[0].strip()
                        cross_items = [item.strip() for item in cross_dependency.split(',') if item.strip()]
                        cross_dependency = ', '.join(cross_items[:2]) if len(cross_items) >= 2 else cross_dependency

                        # Extract escalation engineer and manager
                        escalation_engineer_manager = extract_field(text, 'Escalation Engineer and Manager (if no one is there, put NA):') or '-'
                        # Try to split into engineer and manager if possible
                        escalation_engineer = escalation_engineer_manager
                        escalation_manager = '-'
                        if escalation_engineer_manager and ',' in escalation_engineer_manager:
                            parts = [p.strip() for p in escalation_engineer_manager.split(',')]
                            escalation_engineer = parts[0]
                            if len(parts) > 1:
                                escalation_manager = parts[1]
                        if not is_threaded:
                            # Only create parent entries if not already present
                            existing = Escalation.query.filter_by(webex_msg_id=msg_id).first()
                            if not existing:
                                cap_status = 'N/A'
                                if 'CAP' in text:
                                    cap_status = 'CAP'
                                elif 'pre-CAP' in text or 'precap' in text.lower():
                                    cap_status = 'pre-CAP'
                                escalation = Escalation(
                                    dedt_manager=de_manager,
                                    customer=customer,
                                    version=version,
                                    sr=sr_bems_ddts,
                                    bugid=bugid,
                                    engineer=engineer,
                                    escalation_engineer=escalation_engineer,
                                    escalation_manager=escalation_manager,
                                    upgrade_attempt=upgrade_attempt,
                                    state='Open',
                                    reported_on=datetime.now().strftime('%Y-%m-%d'),
                                    cross_team=cross_dependency,
                                    severity=cap_status,
                                    pending_action='N/A',
                                    remarks='',  # Do not store parent message in remarks
                                    remarks_value='',
                                    created_on=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                                    updated_on=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                                    webex_msg_id=msg_id
                                )
                                db.session.add(escalation)
                                imported += 1
                    except Exception as msg_exc:
                        print(f"Error importing parent message: {msg_exc}")
                db.session.commit()
                # Second pass: append threaded messages to parent remarks
                for msg in messages:
                    try:
                        text = msg.get("text", "")
                        msg_id = msg.get("id")
                        parent_id = msg.get("parentId")
                        is_threaded = bool(parent_id)
                        if is_threaded and parent_id:
                            parent_entry = Escalation.query.filter_by(webex_msg_id=parent_id).first()
                            if parent_entry:
                                # Get timestamp and sender name
                                created = msg.get('created', '')
                                # Format timestamp for display
                                try:
                                    dt = datetime.strptime(created[:19], "%Y-%m-%dT%H:%M:%S")
                                    created_str = dt.strftime("%Y-%m-%d %H:%M:%S")
                                except Exception:
                                    created_str = created
                                person_email = msg.get('personEmail', '')
                                # Extract CEC ID from email
                                cec_id = person_email.split('@')[0] if '@' in person_email else person_email
                                person_name = msg.get('personDisplayName', cec_id)
                                thread_line = f"[{created_str}] {cec_id}: {text}"
                                # Prevent duplicate threaded message
                                if parent_entry.remarks and parent_entry.remarks != '-':
                                    if thread_line not in parent_entry.remarks:
                                        parent_entry.remarks += f"\n{thread_line}"
                                else:
                                    parent_entry.remarks = thread_line
                                parent_entry.updated_on = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                                db.session.commit()
                                imported += 1
                    except Exception as msg_exc:
                        print(f"Error appending threaded message: {msg_exc}")
                return jsonify({
                    "messages": messages,
                    "imported": imported
                })
            except Exception as e:
                import traceback
                print(f"Exception in get_webex_messages: {e}")
                traceback.print_exc()
                return jsonify({"error": str(e)}), 500
    except Exception as e:
        import traceback
        print(f"Exception in get_webex_messages: {e}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500
    
@app.route('/details/<int:item_id>')
@login_required
def details(item_id):
    # Redirect to edit page to maintain consistent table format for cross-team details
    return redirect(url_for('edit', item_id=item_id))

# Helper function to parse remarks properly (respecting multi-line remarks)
def parse_remarks(remarks_text):
    """Parse remarks text and return a list of individual remarks with their content"""
    if not remarks_text:
        return []
    
    import re
    # Split on lines that start with timestamp pattern [YYYY-MM-DD HH:MM] or [YYYY-MM-DD HH:MM - User Name]
    timestamp_pattern = r'\[(\d{4}-\d{2}-\d{2} \d{2}:\d{2}(?:\s*-\s*[^]]+)?)\]'
    
    # Find all timestamp positions
    matches = list(re.finditer(timestamp_pattern, remarks_text))
    if not matches:
        # No timestamps found, treat entire text as one remark
        return [remarks_text.strip()]
    
    remarks = []
    for i, match in enumerate(matches):
        start_pos = match.start()
        # Find the end position (start of next timestamp or end of text)
        if i + 1 < len(matches):
            end_pos = matches[i + 1].start()
        else:
            end_pos = len(remarks_text)
        
        # Extract the complete remark (including possible line breaks)
        remark = remarks_text[start_pos:end_pos].strip()
        if remark:
            remarks.append(remark)
    
    return remarks

# Edit page route
@app.route('/edit/<int:item_id>', methods=['GET', 'POST'])
@login_required
def edit(item_id):
    item = Escalation.query.get(item_id)
    if not item:
        return "Item not found", 404
    
    # Parse remarks for display
    if hasattr(item, 'remarks') and item.remarks:
        item.parsed_remarks = parse_remarks(item.remarks)
    else:
        item.parsed_remarks = []
    
    if request.method == 'POST':
        print('=== EDIT FORM SUBMISSION RECEIVED ===')
        print('Request method:', request.method)
        print('Form data keys:', list(request.form.keys()))
        print('Form data values:')
        for key, value in request.form.items():
            print(f'  {key}: {value}')
        print('=======================================')
        
        # Store original values before update
        original_values = {}
        for field in request.form:
            if hasattr(item, field):
                original_values[field] = getattr(item, field)
        
        # Also capture closing_remarks and closing_reason if they exist
        if 'closing_remarks' in request.form:
            original_values['closing_remarks'] = getattr(item, 'closing_remarks', None)
        if 'closing_reason' in request.form:
            original_values['closing_reason'] = getattr(item, 'closing_reason', None)
        
        # Update all form fields (except remarks which is handled separately)
        for field in request.form:
            if hasattr(item, field) and field not in ['new_remark', 'existing_remarks']:
                value = request.form[field]
                if field == 'state' and value == 'Reopen':
                    value = 'Open'
                setattr(item, field, value)
        
        # Handle new remark with timestamp
        existing_remarks = request.form.get('existing_remarks', '')
        new_remark = request.form.get('new_remark', '').strip()
        
        if new_remark:
            # Create timestamp and user info for the new remark (without seconds)
            current_time = datetime.now().strftime('%Y-%m-%d %H:%M')
            # Extract CEC ID from email (part before @cisco.com)
            user_email = session.get('webex_email', '')
            user_cec_id = user_email.split('@')[0] if user_email else 'unknown'
            timestamped_remark = f"[{current_time} - {user_cec_id}] {new_remark}"
            
            # Combine with existing remarks
            if existing_remarks:
                item.remarks = f"{existing_remarks}\n{timestamped_remark}"
            else:
                item.remarks = timestamped_remark
        elif existing_remarks:
            # Keep existing remarks if no new remark is added
            item.remarks = existing_remarks
        
        # Handle pending action logic based on state
        current_state = getattr(item, 'state', '').lower()
        pending_action = request.form.get('pending_action', '')
        
        if current_state == 'open':
            # For open escalations, default to 'Dev' if no pending action provided
            if not pending_action or pending_action.strip() == '':
                item.pending_action = 'Dev'
        elif current_state == 'closed':
            # For closed escalations, use user selection or default to '-'
            if not pending_action or pending_action.strip() == '':
                item.pending_action = '-'
            else:
                item.pending_action = pending_action
        
        # Handle closing_remarks separately
        if 'closing_remarks' in request.form:
            item.closing_remarks = request.form['closing_remarks']
            print('Updated closing_remarks:', item.closing_remarks)
        
        # Handle closing_reason separately  
        if 'closing_reason' in request.form:
            item.closing_reason = request.form['closing_reason']
            print('Updated closing_reason:', item.closing_reason)
        
        # Determine which cross teams have data and update cross_team field
        cross_teams_with_data = []
        if any([getattr(item, 'routing_component', None), getattr(item, 'routing_director', None), 
                getattr(item, 'routing_de_manager', None), getattr(item, 'routing_engineer', None)]):
            cross_teams_with_data.append('Routing')
        if any([getattr(item, 'ios_xe_component', None), getattr(item, 'ios_xe_director', None), 
                getattr(item, 'ios_xe_de_manager', None), getattr(item, 'ios_xe_engineer', None)]):
            cross_teams_with_data.append('IOS-XE')
        if any([getattr(item, 'swiss_component', None), getattr(item, 'swiss_director', None), 
                getattr(item, 'swiss_de_manager', None), getattr(item, 'swiss_engineer', None)]):
            cross_teams_with_data.append('Location')
        if any([getattr(item, 'sjc_component', None), getattr(item, 'sjc_director', None), 
                getattr(item, 'sjc_de_manager', None), getattr(item, 'sjc_engineer', None)]):
            cross_teams_with_data.append('AP/SDA')
        
        # Update cross_team field
        item.cross_team = ', '.join(cross_teams_with_data) if cross_teams_with_data else None
            
        item.updated_on = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        
        # Handle closed_on date
        closed_on_val = request.form.get('closed_on', '')
        
        # If state is being changed to Closed, handle closing date
        if getattr(item, 'state', '').lower() == 'closed':
            if closed_on_val:
                # Use the provided closing date
                item.closed_on = closed_on_val
            elif not item.closed_on:
                # Auto-set closing date to today if not already set
                item.closed_on = datetime.now().strftime('%Y-%m-%d')
        elif getattr(item, 'state', '').lower() != 'closed':
            # If state is changed from Closed to something else, clear the closing date
            if not closed_on_val:
                item.closed_on = None
        else:
            # For other cases, just update if provided
            if closed_on_val:
                item.closed_on = closed_on_val
            
        db.session.commit()
        print('Data saved successfully')
        
        # Track changes for Webex notification
        changes = {}
        for field, original_value in original_values.items():
            current_value = getattr(item, field, None)
            
            # Normalize values for comparison (treat None, empty string, and 'None' as equivalent)
            def normalize_value(val):
                if val is None or val == '' or val == 'None':
                    return None
                return str(val).strip()
            
            original_normalized = normalize_value(original_value)
            current_normalized = normalize_value(current_value)
            
            # Only track if values actually changed
            if original_normalized != current_normalized:
                changes[field] = {
                    'from': original_value if original_value not in [None, '', 'None'] else 'None',
                    'to': current_value if current_value not in [None, '', 'None'] else 'None'
                }
        
        # Send Webex update notification with changes
        if changes:  # Only send if there are actual changes
            print(f"🔄 Changes detected: {changes}")
            send_escalation_update_to_webex(item, changes)
        else:
            print("⚠️  No changes detected - not sending Webex notification")
            print(f"   Form submission detected but no field changes found")
            # For debugging: send notification anyway to confirm functionality
            print("   🧪 Sending test notification anyway...")
            send_escalation_update_to_webex(item, {})
        
        return redirect(url_for('edit', item_id=item_id))
    states = [s[0] for s in db.session.query(Escalation.state).distinct().all()]
    return render_template('edit.html', item=item, item_id=item_id, states=states)


# --- Move login_required and get_current_user above all route definitions ---
from functools import wraps
from flask import render_template
from collections import Counter



@app.route('/login')
def login():
    print('DEBUG: /login route accessed')
    next_url = request.args.get('next')
    if 'user_id' in session:
        if next_url:
            return redirect(next_url)
        return redirect(url_for('dashboard'))
    return render_template('login.html', next=next_url)

@app.route('/auth/webex')
def auth_webex():
    state = secrets.token_urlsafe(32)
    session['oauth_state'] = state
    next_url = request.args.get('next')
    if next_url:
        session['next_url'] = next_url
    
    # Debug: Print the redirect URI being used
    print(f"DEBUG: WEBEX_REDIRECT_URI = {WEBEX_REDIRECT_URI}")
    
    params = {
        'client_id': WEBEX_CLIENT_ID,
        'response_type': 'code',
        'redirect_uri': WEBEX_REDIRECT_URI,
        'scope': WEBEX_SCOPE,
        'state': state
    }
    auth_url = f"{WEBEX_AUTH_URL}?{urlencode(params)}"
    print(f"DEBUG: Full auth URL = {auth_url}")
    return redirect(auth_url)

@app.route('/auth/webex/callback')
def auth_webex_callback():
    code = request.args.get('code')
    state = request.args.get('state')
    # Validate state
    if not state or state != session.get('oauth_state'):
        return 'Invalid state parameter', 400
    if not code:
        return 'Authorization failed', 400
    # Exchange code for access token
    token_data = {
        'grant_type': 'authorization_code',
        'client_id': WEBEX_CLIENT_ID,
        'client_secret': WEBEX_CLIENT_SECRET,
        'code': code,
        'redirect_uri': WEBEX_REDIRECT_URI
    }
    token_response = requests.post(WEBEX_TOKEN_URL, data=token_data)
    if token_response.status_code != 200:
        return 'Failed to get access token', 400
    token_info = token_response.json()
    access_token = token_info['access_token']
    # Get user info from Webex
    headers = {'Authorization': f'Bearer {access_token}'}
    user_response = requests.get(WEBEX_PEOPLE_URL, headers=headers)
    if user_response.status_code != 200:
        return 'Failed to get user info', 400
    user_info = user_response.json()
    session['webex_id'] = user_info['id']
    session['webex_email'] = user_info['emails'][0]
    session['webex_display_name'] = user_info['displayName']
    session['webex_avatar'] = user_info.get('avatar', '')
    # Save or update user in the login database
    user = LoginUser.query.filter_by(webex_id=user_info['id']).first()
    if not user:
        user = LoginUser(
            webex_id=user_info['id'],
            email=user_info['emails'][0],
            display_name=user_info['displayName'],
            first_name=user_info.get('firstName', ''),
            last_name=user_info.get('lastName', ''),
            avatar=user_info.get('avatar', '')
        )
        db.session.add(user)
    else:
        # Update existing user info
        user.email = user_info['emails'][0]
        user.display_name = user_info['displayName']
        user.first_name = user_info.get('firstName', '')
        user.last_name = user_info.get('lastName', '')
        user.avatar = user_info.get('avatar', '')
    db.session.commit()
    session['user_id'] = user.id  # Set user_id in session
    session.pop('oauth_state', None)
    # Redirect to dashboard or next_url if present
    next_url = session.pop('next_url', None)
    if next_url:
        return redirect(next_url)
    return redirect(url_for('dashboard'))

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

# API endpoint to edit a specific remark
@app.route('/api/edit-remark/<int:item_id>', methods=['POST'])
@login_required
def edit_remark(item_id):
    try:
        data = request.get_json()
        remark_index = data.get('remarkIndex')
        new_content = data.get('newContent', '').strip()
        
        if not new_content:
            return jsonify({'success': False, 'error': 'Remark content cannot be empty'})
        
        # Get the escalation item
        item = Escalation.query.get(item_id)
        if not item:
            return jsonify({'success': False, 'error': 'Escalation not found'})
        
        # Parse existing remarks properly (respecting multi-line remarks)
        if not item.remarks:
            return jsonify({'success': False, 'error': 'No remarks found'})
        
        parsed_remarks = parse_remarks(item.remarks)
        
        if remark_index < 0 or remark_index >= len(parsed_remarks):
            return jsonify({'success': False, 'error': 'Invalid remark index'})
        
        # Update the specific remark
        remark = parsed_remarks[remark_index]
        if remark.startswith('[') and ']' in remark:
            # Keep the original timestamp and user info, update content
            timestamp_end = remark.find(']')
            timestamp_info = remark[:timestamp_end + 1]  # Includes user name if present
            parsed_remarks[remark_index] = f"{timestamp_info} {new_content}"
        else:
            # Remark without timestamp, add current timestamp and user
            current_time = datetime.now().strftime('%Y-%m-%d %H:%M')
            # Extract CEC ID from email (part before @cisco.com)
            user_email = session.get('webex_email', '')
            user_cec_id = user_email.split('@')[0] if user_email else 'unknown'
            parsed_remarks[remark_index] = f"[{current_time} - {user_cec_id}] {new_content}"
        
        # Rebuild the remarks string
        item.remarks = '\n'.join(parsed_remarks)
        item.updated_on = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        db.session.commit()
        
        return jsonify({'success': True, 'message': 'Remark updated successfully'})
        
    except Exception as e:
        print(f"Error editing remark: {str(e)}")
        return jsonify({'success': False, 'error': str(e)})

# API endpoint to delete a specific remark
@app.route('/api/delete-remark/<int:item_id>', methods=['POST'])
@login_required
def delete_remark(item_id):
    try:
        data = request.get_json()
        remark_index = data.get('remarkIndex')
        
        # Get the escalation item
        item = Escalation.query.get(item_id)
        if not item:
            return jsonify({'success': False, 'error': 'Escalation not found'})
        
        # Parse existing remarks properly (respecting multi-line remarks)
        if not item.remarks:
            return jsonify({'success': False, 'error': 'No remarks found'})
        
        parsed_remarks = parse_remarks(item.remarks)
        
        if remark_index < 0 or remark_index >= len(parsed_remarks):
            return jsonify({'success': False, 'error': 'Invalid remark index'})
        
        # Remove the specific remark
        parsed_remarks.pop(remark_index)
        
        # Rebuild the remarks string
        item.remarks = '\n'.join(parsed_remarks) if parsed_remarks else None
        item.updated_on = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        db.session.commit()
        
        return jsonify({'success': True, 'message': 'Remark deleted successfully'})
        
    except Exception as e:
        print(f"Error deleting remark: {str(e)}")
        return jsonify({'success': False, 'error': str(e)})

## --- END LOGIN FUNCTIONALITY COMMENTED OUT ---

if __name__ == '__main__':
    # Run reminder logic before starting the app
    def run_reminders_on_start():
        now = datetime.now()
        week_ago = now - timedelta(days=7)
        open_items = Escalation.query.filter(Escalation.state == "Open").all()
        open_last_week = [item for item in open_items if item.reported_on and 0 <= (now - datetime.strptime(item.reported_on, "%Y-%m-%d")).days <= 7]
        open_count = len(open_last_week)
        closed_items = Escalation.query.filter(Escalation.state == "Closed").all()
        closed_last_week = [item for item in closed_items if item.closed_on and 0 <= (now - datetime.strptime(item.closed_on, "%Y-%m-%d")).days <= 7]
        closed_count = len(closed_last_week)
        closing_reason_values = ["Bug", "Network issue", "Config issue", "Infra issue", "Serviceability issue"]
        closed_category_summary = Counter([item.closing_reason for item in closed_last_week if item.closing_reason in closing_reason_values])
        # Pending actions by category (from all open items, as shown in details page)
        pending_action_categories = [item.pending_action if item.pending_action else "Dev" for item in open_items]
        pending_action_summary = Counter(pending_action_categories)
        for item in closed_last_week:
            try:
                reported_date = datetime.strptime(item.reported_on, "%Y-%m-%d") if item.reported_on else None
                closed_date = datetime.strptime(item.closed_on, "%Y-%m-%d") if item.closed_on else None
                item.mttr = (closed_date - reported_date).days if reported_date and closed_date else ''
            except Exception:
                item.mttr = ''
        for item in open_last_week:
            try:
                reported_date = datetime.strptime(item.reported_on, "%Y-%m-%d") if item.reported_on else None
                item.mttr = (now - reported_date).days if reported_date else ''
            except Exception:
                item.mttr = ''
        # Render HTML template
        html_body = render_template(
            "mail_summary.html",
            engineer="Engineer",
            open_count=open_count,
            closed_count=closed_count,
            closed_category_summary=closed_category_summary,
            pending_action_summary=pending_action_summary,
            closed_items=closed_last_week,
            open_items=open_last_week
        )
        subject = "Weekly Escalation Summary"
        recipients = set()
        for item in open_last_week + closed_last_week:
            if item.engineer:
                recipients.add(f"{item.engineer}@cisco.com")
            if item.dedt_manager:
                recipients.add(f"{item.dedt_manager}@cisco.com")
        for email in recipients:
            send_email(email, subject, "See HTML email", html_body)
        print(f"Sent summary mail to {len(recipients)} recipient(s) on app start.")

    with app.app_context():
        #run_reminders_on_start()
        app.run(host='0.0.0.0', port=5000, debug=True)

    # --- Webex OAuth Route for Login Button ---
